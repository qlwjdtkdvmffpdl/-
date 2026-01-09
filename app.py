import os
import re
import json
import math
import base64
from io import BytesIO
from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple, Any

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib
from matplotlib import font_manager

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm
from reportlab.lib.utils import ImageReader
from reportlab.platypus import Table, TableStyle
from reportlab.lib import colors

from pypdf import PdfReader
from pptx import Presentation
from docx import Document

from langchain_openai import ChatOpenAI
from langchain_core.messages import SystemMessage, HumanMessage, AIMessage

# =========================================================
# Matplotlib font setup (Korean-safe) for charts
# =========================================================
def _set_korean_matplotlib_font() -> None:
    """Try to set a Korean-capable font if available to avoid garbled labels."""
    try:
        candidates = [
            "Malgun Gothic",        # Windows
            "AppleGothic",          # macOS
            "NanumGothic",          # Linux (common)
            "Noto Sans CJK KR",     # Google Noto CJK
            "Noto Sans KR",
            "NanumBarunGothic",
        ]
        available = {f.name for f in font_manager.fontManager.ttflist}
        for c in candidates:
            if c in available:
                matplotlib.rcParams["font.family"] = c
                break
        # prevent minus sign from rendering as a square
        matplotlib.rcParams["axes.unicode_minus"] = False
    except Exception:
        # If font probing fails, keep defaults
        pass

# Apply once at import time (safe to call multiple times).
_set_korean_matplotlib_font()


# =========================================================
# Product/PoV defaults (no model name in UI, fixed internally)
# =========================================================
MODEL_NAME = "gpt-4o-mini"

# Keep PoV stable
MAX_CONTEXT_CHARS = 14000
MAX_IMAGES = 3
MAX_ANALYSIS_ROWS = 200000   # hard cap for safety
MAX_PREVIEW_ROWS = 200

# =========================================================
# Utilities
# =========================================================
STOPWORDS = {
    "a","an","the","and","or","but","to","of","in","on","for","with",
    "is","are","was","were","be","as","at","by","it","this","that",
    "저","그","이","것","수","등","및","또는","그리고","에서","으로",
    "입니다","합니다","관련","대해","기준","요청","확인","정리"
}

def _safe_text(x) -> str:
    return x if isinstance(x, str) else ""

def _truncate(s: str, max_chars: int = 4000) -> str:
    if len(s) <= max_chars:
        return s
    return s[: max_chars - 50] + "\n...(truncated)...\n"

def _keywords(query: str, max_k: int = 10) -> List[str]:
    raw = re.split(r"[\s,./(){}\[\]:;\"'<>!?|\\]+", (query or "").strip().lower())
    ks = []
    for w in raw:
        if len(w) < 2:
            continue
        if w in STOPWORDS:
            continue
        ks.append(w)
    # de-dup preserve order
    out = []
    for w in ks:
        if w not in out:
            out.append(w)
    return out[:max_k]


def _is_numeric_query(q: str) -> bool:
    """Heuristic: if question requests totals/aggregations/comparisons, route to pandas engine."""
    q = (q or "").lower()
    # Korean + English aggregation keywords
    keys = [
        "총액","총합","합계","합산","합","평균","중앙값","최대","최소","분산","표준편차",
        "비율","퍼센트","%","증가","감소","차이","비교","top","상위","하위","순위",
        "그래프","차트","추세","월별","일별","주별","group","sum","total","average","mean",
        "count","min","max","difference","delta","ratio","percent","trend","chart","plot"
    ]
    return any(k in q for k in keys)

def _is_report_query(q: str) -> bool:
    """If user asks for a report/summary, treat as pandas-first (no LLM arithmetic)."""
    q = (q or "").lower()
    keys = [
        "보고서","리포트","요약","정리","pdf","출력","다운로드",
        "재고보고서","재고 보고서","inventory report","report"
    ]
    return any(k in q for k in keys)

def _coerce_numeric_series(ser: pd.Series) -> pd.Series:
    """Robust numeric coercion for Excel-like strings (commas, currency, %, spaces, NBSP)."""
    if ser is None:
        return ser
    s = ser.astype(str)
    # normalize weird spaces
    s = s.str.replace(r"[\u00A0\u2007\u202F]", " ", regex=True)
    # parentheses negative: (1,234) -> -1234
    s = s.str.replace(r"^\((.*)\)$", r"-\1", regex=True)
    # remove common tokens
    s = s.str.replace(",", "", regex=False)
    s = s.str.replace(r"\s+", "", regex=True)
    s = s.str.replace(r"(원|₩|￦|%|퍼센트)", "", regex=True)
    return pd.to_numeric(s, errors="coerce")

def _pick_col(df: pd.DataFrame, candidates: List[str]) -> Optional[str]:
    for c in candidates:
        if c in df.columns:
            return c
    return None

def build_inventory_report_markdown(df: pd.DataFrame, dataset_file: str) -> str:
    """Deterministic inventory report. Numbers MUST come from pandas only."""
    provenance = f"[근거: {dataset_file}]"
    if df is None or df.empty:
        return f"데이터가 비어 있습니다.\n\n{provenance}"

    # Column mapping (best-effort, deterministic)
    col_item = _pick_col(df, ["상품명","품목명","item","name"])
    col_code = _pick_col(df, ["품목코드","상품코드","code","sku"])
    col_qty_total = _pick_col(df, ["전산수량","전산 수량","수량","qty","quantity"])
    col_loc = _pick_col(df, ["보관 위치","보관위치","로케이션","location"])
    col_obsolete_qty = _pick_col(df, ["불용재고","불용","폐기","obsolete_qty","obsolete"])
    col_normal_qty = _pick_col(df, ["정상품","정상재고","정상","normal_qty","normal"])
    col_unit_price = _pick_col(df, ["제품단가","단가","unit_price","price"])
    col_obsolete_amt = _pick_col(df, ["불용재고 총액","불용재고총액","불용금액","obsolete_amount","불용총액"])
    col_normal_amt = _pick_col(df, ["정상재고 총액","정상재고총액","정상금액","normal_amount","정상총액"])

    work = df.copy()

    # Ensure numeric columns are numeric
    for c in [col_obsolete_qty, col_normal_qty, col_unit_price, col_obsolete_amt, col_normal_amt]:
        if c and c in work.columns:
            work[c] = _coerce_numeric_series(work[c])

    # Derive amounts if missing
    if (not col_obsolete_amt) and col_obsolete_qty and col_unit_price:
        col_obsolete_amt = "불용재고 총액"
        work[col_obsolete_amt] = work[col_obsolete_qty] * work[col_unit_price]
    if (not col_normal_amt) and col_normal_qty and col_unit_price:
        col_normal_amt = "정상재고 총액"
        work[col_normal_amt] = work[col_normal_qty] * work[col_unit_price]

    # Compute totals (skipna True)
    n_items = int(len(work))
    total_obsolete_qty = int(_coerce_numeric_series(work[col_obsolete_qty]).sum()) if col_obsolete_qty else None
    total_normal_qty = int(_coerce_numeric_series(work[col_normal_qty]).sum()) if col_normal_qty else None
    total_obsolete_amt = float(_coerce_numeric_series(work[col_obsolete_amt]).sum()) if col_obsolete_amt else None
    total_normal_amt = float(_coerce_numeric_series(work[col_normal_amt]).sum()) if col_normal_amt else None

    # Build a display table with preferred columns (only those that exist)
    preferred = []
    for c in [col_code, col_item, col_qty_total, col_loc, col_obsolete_qty, col_normal_qty, col_unit_price, col_obsolete_amt, col_normal_amt]:
        if c and c in work.columns and c not in preferred:
            preferred.append(c)
    table_df = work[preferred] if preferred else work

    # Render markdown (numbers are injected from pandas)
    md_lines = []
    md_lines.append("# 재고보고서")
    md_lines.append("\n## 1. 기본 정보")
    md_lines.append(f"- 총 품목 수: **{n_items:,}개**")
    md_lines.append("\n## 2. 재고 현황")
    md_lines.append(table_df.head(50).to_markdown(index=False))
    md_lines.append("\n## 3. 요약")
    if total_obsolete_qty is not None:
        md_lines.append(f"- 총 불용재고 수량: **{total_obsolete_qty:,}**")
    if total_normal_qty is not None:
        md_lines.append(f"- 총 정상재고 수량: **{total_normal_qty:,}**")
    if total_obsolete_amt is not None and not math.isnan(total_obsolete_amt):
        md_lines.append(f"- 총 불용재고 총액: **{int(round(total_obsolete_amt)):,}원**")
    if total_normal_amt is not None and not math.isnan(total_normal_amt):
        md_lines.append(f"- 총 정상재고 총액: **{int(round(total_normal_amt)):,}원**")
    md_lines.append("\n## 다음 액션 제안")
    md_lines.append("1. 불용재고 처리 방안 검토\n2. 정상재고 판매 계획 수립")
    md_lines.append(f"\n{provenance}")
    return "\n".join(md_lines)

def _needs_visualization(q: str) -> bool:
    """
    Determine if the query needs table/chart visualization or just a simple text answer.
    Returns:
        True: needs table/chart (comparison, breakdown, ranking, etc.)
        False: simple numeric answer is enough (e.g., "What's the total?")
    """
    q = (q or "").lower()

    # Explicit visualization requests
    viz_keywords = [
        "그래프", "차트", "chart", "plot", "시각화", "비주얼",
        "표", "테이블", "table", "목록", "리스트", "list"
    ]
    if any(k in q for k in viz_keywords):
        return True

    # Comparison/breakdown/analysis (usually needs table)
    analysis_keywords = [
        "비교", "별로", "별", "각", "top", "상위", "하위", "순위", "ranking",
        "분석", "분포", "breakdown", "group", "나눠", "나누어", "구분"
    ]
    if any(k in q for k in analysis_keywords):
        return True

    # Simple questions (just need a number/text answer)
    simple_keywords = ["얼마", "몇", "언제", "어디", "무엇", "누구"]
    # If it's a simple question WITHOUT comparison/breakdown keywords
    if any(k in q for k in simple_keywords):
        # Double-check it's not a complex question
        if not any(k in q for k in analysis_keywords):
            return False

    # Default: safe to generate visualization (conservative approach)
    return True

def _format_result_markdown(result_df: pd.DataFrame, provenance: str) -> str:
    """Create a deterministic response. No LLM arithmetic."""
    if result_df is None or result_df.empty:
        return f"결과가 비어 있습니다.\n\n{provenance}"
    # If single-row aggregates, print key-values
    if len(result_df) == 1 and len(result_df.columns) <= 8:
        row = result_df.iloc[0].to_dict()
        lines = []
        for k, v in row.items():
            # pretty number
            if isinstance(v, (int, float)) and not pd.isna(v):
                vv = int(v) if float(v).is_integer() else float(v)
                lines.append(f"- **{k}**: {vv:,}")
            else:
                lines.append(f"- **{k}**: {v}")
        return "\n".join(lines) + f"\n\n{provenance}"
    # Otherwise show table note (table will be rendered separately)
    return f"요청하신 집계/비교 결과를 표로 생성했습니다.\n\n{provenance}"

def _infer_numeric_columns(df: pd.DataFrame) -> List[str]:
    cols = []
    for c in df.columns:
        try:
            s = pd.to_numeric(df[c], errors="coerce")
            if s.notna().mean() >= 0.6:
                cols.append(str(c))
        except Exception:
            continue
    return cols

def _infer_datetime_columns(df: pd.DataFrame) -> List[str]:
    cols = []
    for c in df.columns:
        try:
            s = pd.to_datetime(df[c], errors="coerce", infer_datetime_format=True)
            if s.notna().mean() >= 0.6:
                cols.append(str(c))
        except Exception:
            continue
    return cols

def _df_relevant_rows(df: pd.DataFrame, query: str, max_rows: int = MAX_PREVIEW_ROWS) -> pd.DataFrame:
    ks = _keywords(query)
    if df is None or df.empty:
        return df
    if not ks:
        return df.head(min(max_rows, len(df)))

    # cap for stability
    sample_df = df
    if len(df) > MAX_ANALYSIS_ROWS:
        sample_df = df.head(MAX_ANALYSIS_ROWS)

    s = sample_df.astype(str).apply(lambda col: col.str.lower())
    mask = None
    for k in ks:
        col_mask = s.apply(lambda col: col.str.contains(re.escape(k), na=False))
        row_mask = col_mask.any(axis=1)
        mask = row_mask if mask is None else (mask | row_mask)

    filtered = sample_df[mask] if mask is not None else sample_df
    if filtered.empty:
        return sample_df.head(min(max_rows, len(sample_df)))
    return filtered.head(min(max_rows, len(filtered)))

def _apply_filters(df: pd.DataFrame, filters: List[Dict[str, Any]]) -> pd.DataFrame:
    """Filters schema:
    [{"column":"col","op":"eq|neq|contains|gt|gte|lt|lte|in|between","value":...}, ...]
    """
    out = df
    for f in filters or []:
        col = f.get("column")
        op = f.get("op")
        val = f.get("value")
        if col not in out.columns:
            continue
        try:
            if op == "eq":
                out = out[out[col] == val]
            elif op == "neq":
                out = out[out[col] != val]
            elif op == "contains":
                out = out[out[col].astype(str).str.contains(str(val), case=False, na=False)]
            elif op in ("gt","gte","lt","lte"):
                s = pd.to_numeric(out[col], errors="coerce")
                v = float(val)
                if op == "gt":
                    out = out[s > v]
                elif op == "gte":
                    out = out[s >= v]
                elif op == "lt":
                    out = out[s < v]
                elif op == "lte":
                    out = out[s <= v]
            elif op == "in":
                if isinstance(val, list):
                    out = out[out[col].isin(val)]
            elif op == "between":
                if isinstance(val, list) and len(val) == 2:
                    s = pd.to_numeric(out[col], errors="coerce")
                    lo, hi = float(val[0]), float(val[1])
                    out = out[(s >= lo) & (s <= hi)]
        except Exception:
            # keep stable: ignore bad filters
            continue
    return out

# =========================================================
# File cache schema
# =========================================================
@dataclass
class CachedFile:
    name: str
    ext: str
    kind: str  # "table" | "pdf" | "ppt" | "docx" | "image" | "unknown"
    meta: Dict
    df: Optional[pd.DataFrame] = None
    pdf_pages: Optional[List[str]] = None
    ppt_slides: Optional[List[Dict]] = None
    docx_text: Optional[str] = None
    image_b64: Optional[str] = None

def parse_uploaded_file(uploaded) -> CachedFile:
    name = uploaded.name
    ext = name.split(".")[-1].lower()

    if ext in ["xlsx", "csv"]:
        try:
            if ext == "xlsx":
                df = pd.read_excel(uploaded)
            else:
                df = pd.read_csv(uploaded)
        except Exception as e:
            return CachedFile(name=name, ext=ext, kind="unknown", meta={"error": str(e)})

        meta = {
            "shape": [int(df.shape[0]), int(df.shape[1])],
            "columns": list(map(str, df.columns.tolist())),
        }
        return CachedFile(name=name, ext=ext, kind="table", meta=meta, df=df)

    if ext == "pdf":
        pages: List[str] = []
        try:
            reader = PdfReader(uploaded)
            for p in reader.pages:
                txt = _safe_text(p.extract_text())
                pages.append(txt)
        except Exception as e:
            return CachedFile(name=name, ext=ext, kind="pdf", meta={"error": str(e)}, pdf_pages=[])

        meta = {"pages": len(pages)}
        empty_pages = sum(1 for t in pages if len(t.strip()) == 0)
        meta["empty_pages"] = empty_pages
        meta["likely_scanned"] = bool(len(pages) > 0 and empty_pages / max(1, len(pages)) >= 0.6)
        return CachedFile(name=name, ext=ext, kind="pdf", meta=meta, pdf_pages=pages)

    if ext == "pptx":
        slides: List[Dict] = []
        try:
            prs = Presentation(uploaded)
            for i, slide in enumerate(prs.slides, start=1):
                texts = []
                for shp in slide.shapes:
                    if hasattr(shp, "text"):
                        t = _safe_text(shp.text).strip()
                        if t:
                            texts.append(t)
                title = texts[0] if texts else f"Slide {i}"
                slides.append({"slide": i, "title": _truncate(title, 120), "text": "\n".join(texts)})
        except Exception as e:
            return CachedFile(name=name, ext=ext, kind="ppt", meta={"error": str(e)}, ppt_slides=[])

        meta = {"slides": len(slides)}
        return CachedFile(name=name, ext=ext, kind="ppt", meta=meta, ppt_slides=slides)

    if ext == "docx":
        try:
            doc = Document(uploaded)
            paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
            text = "\n".join(paras)
        except Exception as e:
            return CachedFile(name=name, ext=ext, kind="docx", meta={"error": str(e)}, docx_text="")

        meta = {"chars": len(text)}
        return CachedFile(name=name, ext=ext, kind="docx", meta=meta, docx_text=text)

    if ext in ["png", "jpg", "jpeg", "webp"]:
        try:
            data = uploaded.getvalue()
            b64 = base64.b64encode(data).decode("utf-8")
        except Exception as e:
            return CachedFile(name=name, ext=ext, kind="image", meta={"error": str(e)}, image_b64=None)

        meta = {"bytes": len(uploaded.getvalue())}
        return CachedFile(name=name, ext=ext, kind="image", meta=meta, image_b64=b64)

    return CachedFile(name=name, ext=ext, kind="unknown", meta={"note": "unsupported file type"})

# =========================================================
# ArtifactStore: 멀티포맷 통합 저장소 (표준화된 아티팩트 저장)
# =========================================================
@dataclass
class Artifact:
    """표준화된 파일 아티팩트. 모든 파일 타입을 동일한 구조로 저장."""
    artifact_id: str
    file_id: str
    file_name: str
    file_type: str  # "excel" | "csv" | "pdf" | "ppt" | "docx" | "image" | "unknown"
    source_loc: str  # excel: "sheet:Sheet1/table:Table1", pdf: "page:3", ppt: "slide:5", image: "file"
    text: str  # 검색/근거용 텍스트
    tables: List[pd.DataFrame]  # 숫자 계산용 DataFrame 리스트
    provenance: str  # 출처 추적용 문자열

class ArtifactStore:
    """멀티포맷 파일을 표준화된 Artifact로 저장하고 조회하는 저장소."""
    
    def __init__(self):
        self.artifacts: Dict[str, Artifact] = {}
    
    def add_from_cached_file(self, cached: CachedFile) -> List[str]:
        """CachedFile을 Artifact로 변환하여 저장. 반환: artifact_id 리스트."""
        artifact_ids = []
        
        if cached.kind == "table" and cached.df is not None:
            # Excel/CSV: DataFrame을 테이블로 저장
            artifact_id = f"table_{cached.name}"
            artifact = Artifact(
                artifact_id=artifact_id,
                file_id=cached.name,
                file_name=cached.name,
                file_type="excel" if cached.ext == "xlsx" else "csv",
                source_loc=f"sheet:default/table:main",
                text="",  # 테이블은 text보다 tables가 중요
                tables=[cached.df],
                provenance=f"[근거: {cached.name}]"
            )
            self.artifacts[artifact_id] = artifact
            artifact_ids.append(artifact_id)
        
        elif cached.kind == "pdf" and cached.pdf_pages is not None:
            # PDF: 각 페이지를 별도 아티팩트로 저장
            for i, page_text in enumerate(cached.pdf_pages, start=1):
                artifact_id = f"pdf_{cached.name}_p{i}"
                artifact = Artifact(
                    artifact_id=artifact_id,
                    file_id=cached.name,
                    file_name=cached.name,
                    file_type="pdf",
                    source_loc=f"page:{i}",
                    text=page_text,
                    tables=[],  # PDF에서 테이블 추출은 향후 확장 가능
                    provenance=f"[근거: {cached.name} / p.{i}]"
                )
                self.artifacts[artifact_id] = artifact
                artifact_ids.append(artifact_id)
        
        elif cached.kind == "ppt" and cached.ppt_slides is not None:
            # PPT: 각 슬라이드를 별도 아티팩트로 저장
            for slide in cached.ppt_slides:
                slide_num = slide.get("slide", 0)
                artifact_id = f"ppt_{cached.name}_s{slide_num}"
                slide_text = f"{slide.get('title', '')}\n{slide.get('text', '')}"
                artifact = Artifact(
                    artifact_id=artifact_id,
                    file_id=cached.name,
                    file_name=cached.name,
                    file_type="ppt",
                    source_loc=f"slide:{slide_num}",
                    text=slide_text,
                    tables=[],
                    provenance=f"[근거: {cached.name} / slide {slide_num}]"
                )
                self.artifacts[artifact_id] = artifact
                artifact_ids.append(artifact_id)
        
        elif cached.kind == "docx" and cached.docx_text:
            # DOCX: 전체 문서를 하나의 아티팩트로
            artifact_id = f"docx_{cached.name}"
            artifact = Artifact(
                artifact_id=artifact_id,
                file_id=cached.name,
                file_name=cached.name,
                file_type="docx",
                source_loc="file",
                text=cached.docx_text,
                tables=[],
                provenance=f"[근거: {cached.name}]"
            )
            self.artifacts[artifact_id] = artifact
            artifact_ids.append(artifact_id)
        
        elif cached.kind == "image" and cached.image_b64:
            # 이미지: OCR은 향후 확장, 현재는 메타데이터만
            artifact_id = f"image_{cached.name}"
            artifact = Artifact(
                artifact_id=artifact_id,
                file_id=cached.name,
                file_name=cached.name,
                file_type="image",
                source_loc="file",
                text="",  # OCR 결과는 향후 추가
                tables=[],
                provenance=f"[근거: {cached.name}]"
            )
            self.artifacts[artifact_id] = artifact
            artifact_ids.append(artifact_id)
        
        return artifact_ids
    
    def get_tables(self) -> List[Tuple[str, pd.DataFrame]]:
        """모든 테이블(DataFrame)을 반환. 반환: [(artifact_id, df), ...]"""
        result = []
        for artifact_id, artifact in self.artifacts.items():
            for df in artifact.tables:
                if df is not None and not df.empty:
                    result.append((artifact_id, df))
        return result
    
    def search_text(self, query: str, max_results: int = 10) -> List[Artifact]:
        """텍스트 검색. 반환: 관련 아티팩트 리스트."""
        keywords = _keywords(query)
        if not keywords:
            return []
        
        scored = []
        for artifact in self.artifacts.values():
            if not artifact.text:
                continue
            text_lower = artifact.text.lower()
            score = sum(1 for kw in keywords if kw in text_lower)
            if score > 0:
                scored.append((score, artifact))
        
        scored.sort(reverse=True, key=lambda x: x[0])
        return [a for _, a in scored[:max_results]]

# =========================================================
# Context builders (Q&A mode)
# =========================================================
def build_context_blocks(files: List[CachedFile], user_query: str) -> Tuple[str, List[Dict]]:
    ks = _keywords(user_query)
    context_parts: List[str] = []
    images: List[Dict] = []

    for f in files:
        if f.kind == "table" and f.df is not None:
            df = f.df
            shape = f.meta.get("shape", ["?", "?"])
            cols = f.meta.get("columns", [])
            sample = _df_relevant_rows(df, user_query, max_rows=200)

            part = [
                f"[TABLE] {f.name}",
                f"- shape: {shape[0]} rows x {shape[1]} cols",
                f"- columns: {', '.join(map(str, cols[:60]))}{' ...' if len(cols) > 60 else ''}",
                "- relevant rows (top):",
                sample.head(min(50, len(sample))).to_markdown(index=False),
            ]
            context_parts.append("\n".join(part))

        elif f.kind == "pdf" and f.pdf_pages is not None:
            pages = f.pdf_pages
            total = len(pages)
            if total == 0:
                context_parts.append(f"[PDF] {f.name}\n- (no extractable text)")
                continue

            picked: List[int] = []
            if ks:
                for i, txt in enumerate(pages, start=1):
                    low = txt.lower()
                    if any(k in low for k in ks):
                        picked.append(i)
                    if len(picked) >= 4:
                        break
            if not picked:
                picked = [1, 2] if total >= 2 else [1]

            chunk = [f"[PDF] {f.name} (pages: {total})"]
            if f.meta.get("likely_scanned"):
                chunk.append("- note: many pages have no text (likely scanned PDF). Answers may be limited without OCR.")
            for pno in picked:
                txt = _truncate(_safe_text(pages[pno - 1]).strip(), 2500)
                if not txt.strip():
                    txt = "(no extractable text on this page)"
                chunk.append(f"- p.{pno}:\n{txt}")
            context_parts.append("\n".join(chunk))

        elif f.kind == "ppt" and f.ppt_slides is not None:
            slides = f.ppt_slides
            total = len(slides)
            if total == 0:
                context_parts.append(f"[PPT] {f.name}\n- (no extractable text)")
                continue

            picked_slides: List[Dict] = []
            if ks:
                for s in slides:
                    low = (s.get("title","") + "\n" + s.get("text","")).lower()
                    if any(k in low for k in ks):
                        picked_slides.append(s)
                    if len(picked_slides) >= 4:
                        break
            if not picked_slides:
                picked_slides = slides[:2]

            chunk = [f"[PPT] {f.name} (slides: {total})"]
            for s in picked_slides:
                sn = s.get("slide", "?")
                title = s.get("title", "")
                txt = _truncate(_safe_text(s.get("text", "")).strip(), 2000)
                chunk.append(f"- slide {sn} ({title}):\n{txt}")
            context_parts.append("\n".join(chunk))

        elif f.kind == "docx":
            txt = _truncate(_safe_text(f.docx_text).strip(), 4000)
            context_parts.append(f"[DOCX] {f.name}\n{txt}")

        elif f.kind == "image" and f.image_b64:
            context_parts.append(f"[IMAGE] {f.name} (provided)")
            if len(images) < MAX_IMAGES:
                images.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/{f.ext};base64,{f.image_b64}"}
                })
        else:
            context_parts.append(f"[FILE] {f.name}\n- note: unsupported or failed to parse. meta={f.meta}")

        if sum(len(p) for p in context_parts) > MAX_CONTEXT_CHARS:
            context_parts.append("\n[NOTE] Context truncated for stability. Ask for a specific file/page/sheet if needed.\n")
            break

    context_text = _truncate("\n\n".join(context_parts), MAX_CONTEXT_CHARS)
    return context_text, images

# =========================================================
# System prompts
# =========================================================
def build_system_prompt(persona: str, mode: str) -> str:
    base = """당신은 기업 실무(물류/자재 정산, 보고, 근거 정리)를 돕는 AI 어시스턴트입니다.

[핵심 규칙]
1) 답변은 반드시 '근거(출처)'를 포함하세요.
   - 근거 표기 형식 예:
     • [근거: 파일명 / sheet:시트명 / rows:행범위]
     • [근거: 파일명 / p.N]
     • [근거: 파일명 / slide N]
     • [근거: 이미지 파일명]
2) 근거가 부족하면 '근거 부족'이라고 명확히 말하고, 어떤 자료가 더 필요할지 제안하세요.
3) 불확실한 추측은 금지합니다. 확실한 것/불확실한 것을 분리하세요.
"""

    if mode == "analysis":
        base += """
[분석 모드 규칙]
- 당신은 '분석 계획(표/그래프/집계)'을 먼저 제안할 수 있습니다.
- 숫자 계산/집계는 파이썬(pandas)이 수행합니다. 당신은 무엇을 계산할지 '계획'을 JSON으로 제시해야 합니다.
- JSON 외의 텍스트는 출력하지 마세요. (파싱 실패 방지)
"""

    personas = {
        "물류 비서": "말투는 간결하고 실무적으로. 가능한 표/리스트로 정리하고, 다음 액션을 제안하세요.",
        "감사관(정산 검증)": "잠재 오류/누락/불일치를 우선 탐지하세요. 체크리스트와 리스크를 제시하고, '추가 확인 필요'를 명확히 표시하세요.",
        "기획자(PPT 요약)": "핵심 메시지/숫자/근거를 요약하고, 슬라이드 구조 제안을 포함하세요."
    }
    return base + "\n[페르소나 지침]\n" + personas.get(persona, personas["물류 비서"])


# =========================================================
# Rule extraction (docs -> JSON rules) and application (pandas)
# =========================================================
RULES_SCHEMA = {
  "source_files": ["doc1.pdf", "doc2.pptx"],
  "date_range": {"column": "date_col_optional", "start": "YYYY-MM-DD", "end": "YYYY-MM-DD"},
  "filters": [
    {"column": "col", "op": "eq|neq|contains|in|gt|gte|lt|lte", "value": "..." }
  ],
  "derived_columns": [
    {"name": "new_col", "formula": "qty * unit_price", "if_missing": True}
  ],
  "overrides": [
    {"target_column": "unit_price", "value": 1234, "when": {"column":"item","op":"eq|contains|in","value":"..."}}
  ],
  "notes": "optional"
}

def _doc_text_for_rules(files: List["CachedFile"], max_chars_per_file: int = 6000) -> str:
    parts = []
    for f in files:
        if f.kind == "pdf" and f.pdf_pages is not None:
            # take first few pages with text
            txt = "\n".join((t or "") for t in f.pdf_pages[:6])
            parts.append(f"[PDF] {f.name}\n{_truncate(txt, max_chars_per_file)}")
        elif f.kind == "ppt" and f.ppt_slides is not None:
            # take first few slides
            txt = "\n".join(f"slide {s.get('slide')}: {s.get('text','')}" for s in f.ppt_slides[:8])
            parts.append(f"[PPT] {f.name}\n{_truncate(txt, max_chars_per_file)}")
        elif f.kind == "docx" and f.docx_text is not None:
            parts.append(f"[DOCX] {f.name}\n{_truncate(f.docx_text, max_chars_per_file)}")
    return "\n\n".join(parts)

def _rules_prompt(df: pd.DataFrame, docs_text: str, doc_names: List[str]) -> str:
    cols = list(map(str, df.columns.tolist()))
    num_cols = _infer_numeric_columns(df)[:25]
    dt_cols = _infer_datetime_columns(df)[:10]
    sample = df.head(20).to_dict(orient="records")

    return f"""
당신은 '문서 규칙(단가/기간/조건)'을 데이터(엑셀/CSV)에 적용하기 위한 규칙 추출기입니다.
아래 문서 내용에서, 데이터 분석/정산에 영향을 주는 '기간', '조건(필터)', '단가/패널티(값/계산식)', '파생 컬럼'을 찾아 JSON 규칙으로만 출력하세요.
절대 다른 텍스트를 출력하지 마세요. JSON만.

[규칙 스키마]
{json.dumps(RULES_SCHEMA, ensure_ascii=False, indent=2)}

[문서 파일명]
{doc_names}

[데이터 파일 컬럼]
- 전체 컬럼: {cols}
- 수치형 후보: {num_cols}
- 날짜형 후보: {dt_cols}

[데이터 샘플(상위 20행)]
{json.dumps(sample, ensure_ascii=False, indent=2)}

[문서 텍스트 발췌]
{docs_text}

[출력 지침]
- 문서에 명시된 기간이 있으면 date_range로 표현하되, 가능한 경우 데이터의 날짜 컬럼명(column)을 함께 지정하세요.
- 문서에 특정 조건(예: 특정 거래처/상태/구간)이 있으면 filters로 표현하세요.
- 문서에 단가/금액 계산 규칙이 있으면 derived_columns 또는 overrides로 표현하세요.
- 확실하지 않으면 비워두고 notes에 '근거 부족'을 적으세요.
"""

def _safe_eval_formula(df: pd.DataFrame, formula: str) -> pd.Series:
    """
    Very small safe expression evaluator for formulas like: qty * unit_price
    Allowed: column names (letters/digits/underscore), numbers, + - * / ( ) and spaces.
    """
    if not isinstance(formula, str):
        raise ValueError("formula must be str")
    if len(formula) > 200:
        raise ValueError("formula too long")

    if re.search(r"[^0-9a-zA-Z_\s\+\-\*\/\(\)\.]", formula):
        raise ValueError("formula contains unsupported characters")

    # Ensure referenced names exist
    tokens = re.findall(r"[A-Za-z_][A-Za-z0-9_]*", formula)
    for t in set(tokens):
        if t not in df.columns:
            raise ValueError(f"Unknown column in formula: {t}")

    # pandas.eval with engine='python' to avoid numexpr surprises
    return pd.eval(formula, engine="python", local_dict={c: df[c] for c in df.columns})

def apply_rules_to_df(df: pd.DataFrame, rules: dict) -> Tuple[pd.DataFrame, str]:
    """
    Apply rules to df safely.
    Returns modified df and a provenance string describing applied rules.
    """
    if not isinstance(rules, dict):
        return df, ""

    work = df.copy()
    prov_parts = []

    # Date range filter
    dr = rules.get("date_range") or {}
    dr_col = dr.get("column")
    dr_start = dr.get("start")
    dr_end = dr.get("end")
    if dr_col and dr_col in work.columns and (dr_start or dr_end):
        try:
            col_dt = pd.to_datetime(work[dr_col], errors="coerce")
            mask = pd.Series(True, index=work.index)
            if dr_start:
                mask &= col_dt >= pd.to_datetime(dr_start, errors="coerce")
            if dr_end:
                mask &= col_dt <= pd.to_datetime(dr_end, errors="coerce")
            work = work[mask]
            prov_parts.append(f"기간 필터 적용: {dr_col} in [{dr_start or '-'}, {dr_end or '-'}]")
        except Exception:
            prov_parts.append(f"기간 필터(실패): {dr_col}")

    # Filters
    for flt in (rules.get("filters") or []):
        try:
            col = flt.get("column")
            op = (flt.get("op") or "").lower()
            val = flt.get("value")
            if col not in work.columns:
                continue
            ser = work[col]
            if op == "eq":
                work = work[ser == val]
            elif op == "neq":
                work = work[ser != val]
            elif op == "contains":
                work = work[ser.astype(str).str.contains(str(val), na=False)]
            elif op == "in":
                if isinstance(val, list):
                    work = work[ser.isin(val)]
            elif op in ("gt","gte","lt","lte"):
                num = pd.to_numeric(ser, errors="coerce")
                vnum = float(val) if val is not None else None
                if vnum is None:
                    continue
                if op == "gt":
                    work = work[num > vnum]
                elif op == "gte":
                    work = work[num >= vnum]
                elif op == "lt":
                    work = work[num < vnum]
                else:
                    work = work[num <= vnum]
            prov_parts.append(f"조건 필터 적용: {col} {op} {val}")
        except Exception:
            continue

    # Overrides
    for ov in (rules.get("overrides") or []):
        try:
            tgt = ov.get("target_column")
            value = ov.get("value")
            when = ov.get("when") or {}
            wcol = when.get("column")
            wop = (when.get("op") or "").lower()
            wval = when.get("value")
            if tgt not in work.columns or wcol not in work.columns:
                continue
            mask = pd.Series(False, index=work.index)
            ser = work[wcol]
            if wop == "eq":
                mask = ser == wval
            elif wop == "contains":
                mask = ser.astype(str).str.contains(str(wval), na=False)
            elif wop == "in" and isinstance(wval, list):
                mask = ser.isin(wval)
            if mask.any():
                work.loc[mask, tgt] = value
                prov_parts.append(f"값 오버라이드: {tgt}={value} when {wcol} {wop} {wval}")
        except Exception:
            continue

    # Derived columns
    for dc in (rules.get("derived_columns") or []):
        try:
            name = dc.get("name")
            formula = dc.get("formula")
            if_missing = bool(dc.get("if_missing", True))
            if not name or not formula:
                continue
            if if_missing and name in work.columns:
                continue
            work[name] = _safe_eval_formula(work, formula)
            prov_parts.append(f"파생 컬럼 생성: {name} = {formula}")
        except Exception as e:
            prov_parts.append(f"파생 컬럼(실패): {dc.get('name')} ({e})")

    notes = rules.get("notes")
    if notes:
        prov_parts.append(f"규칙 노트: {notes}")

    provenance = "; ".join(prov_parts)
    return work, provenance

# =========================================================
# Analysis planner (LLM -> JSON -> execute)
# =========================================================
ANALYSIS_SCHEMA = {
    "task": "table|chart|both",
    "dataset_file": "<excel filename>",
    "select_columns": ["optional", "list"],
    "filters": [
        {"column": "col", "op": "eq|neq|contains|gt|gte|lt|lte|in|between", "value": "any"}
    ],
    "groupby": ["col1", "col2"],
    "metrics": [
        {"column": "numeric_col", "agg": "sum|mean|count|min|max"}
    ],
    "sort": {"by": "metric_name_or_column", "ascending": False},
    "top_n": 20,
    "chart": {
        "type": "bar|line|pie",
        "x": "group_col_or_index",
        "y": "metric_name_or_column",
        "title": "optional",
        "ylabel": "optional",
        "xlabel": "optional"
    },
    "notes": "optional"
}

def _analysis_prompt(df: pd.DataFrame, filename: str, user_query: str) -> str:
    cols = list(map(str, df.columns.tolist()))
    num_cols = _infer_numeric_columns(df)[:25]
    dt_cols = _infer_datetime_columns(df)[:10]
    sample = df.head(20).to_dict(orient="records")

    return f"""
당신은 데이터 분석 계획 생성기입니다.
사용자 질문을 보고, 아래 스키마(JSON)로만 답하세요. 절대 다른 텍스트를 출력하지 마세요.

[스키마]
{json.dumps(ANALYSIS_SCHEMA, ensure_ascii=False, indent=2)}

[데이터셋 정보]
- 파일명: {filename}
- 전체 행 수: {len(df)}개 ⚠️ 모든 행이 유효한 데이터입니다. 임의로 제외하지 마세요!
- 전체 컬럼 수: {len(cols)}개
- 컬럼 목록: {cols}
- 숫자형 컬럼: {num_cols}
- 날짜형 컬럼: {dt_cols}
- 데이터 샘플 (처음 {len(sample)}개 행): {sample}

[사용자 질문]
{user_query}

⚠️ [중요 규칙 - 반드시 지켜야 함!]
1. 필터(filters) 사용 조건:
   - 사용자가 "~만", "~제외", "~인 것", "~보다 큰" 등 명시적 조건을 말했을 때만 사용
   - "전체", "모든", "합계", "총" 같은 단어가 있으면 filters는 빈 배열 []로 설정
   - 예: "불용재고 총액은 얼마?" → filters: [] (조건 없음)
   - 예: "A 거래처의 총액은?" → filters: [{{"column": "거래처", "op": "eq", "value": "A"}}]

2. top_n 사용 조건:
   - 사용자가 "상위 N개", "top N", "N위까지" 등을 명시했을 때만 사용
   - 명시하지 않았으면 top_n 필드를 아예 포함하지 마세요
   - 예: "총액은 얼마?" → top_n 필드 없음
   - 예: "상위 5개 거래처" → top_n: 5

3. groupby 사용 조건:
   - "~별", "각", "항목별" 같은 단어가 있으면 사용
   - 단순 합계/평균 질문이면 groupby는 빈 배열 []로 설정
   - 예: "전체 총액은?" → groupby: []
   - 예: "거래처별 총액은?" → groupby: ["거래처"]

4. 전체 집계 vs 그룹별 집계:
   - "전체", "총", "모든", "합계" → groupby: [], filters: []
   - "~별", "각" → groupby: ["컬럼명"], filters: []

[출력 규칙]
- 반드시 JSON 1개만 출력 (설명 금지)
- metric의 결과 컬럼명은 "{'{agg}({col})'}" 형태로 쓰세요. 예: "sum(amount)"
- chart.y는 metric 결과 컬럼명을 사용하세요
- 위 중요 규칙을 위반하면 계산 결과가 틀립니다!
""".strip()

def _safe_json_load(s: str) -> Optional[dict]:
    try:
        # allow codefence
        s2 = s.strip()
        s2 = re.sub(r"^```json\s*|\s*```$", "", s2, flags=re.IGNORECASE | re.MULTILINE).strip()
        return json.loads(s2)
    except Exception:
        return None

# =========================================================
# Intent Routing: 질문 의도 분류 및 팩 선택
# =========================================================
ROUTE_RESULT_SCHEMA = {
    "action": "EXPLAIN|ANALYZE|COMPARE|REPORT|OPTIMIZE",
    "domain": "logistics|inventory|settlement|sales|general",
    "pack": "PACK_NAME",
    "params": {
        "time_hint": "yesterday|today|date_range|none",
        "file_scope": "active_file|all_files|selected_files",
        "key_hint": ["SKU","품목코드","거래처","상품명"],
        "metric_hint": ["amount","qty","매출","정산금액"],
        "top_n": 10,
        "groupby_hint": ["거래처","센터","카테고리"],
        "compare_target": "yesterday_vs_today|A_vs_B|period_vs_period",
        "output": "chat|pdf"
    },
    "risk": "low|high",
    "confidence": 0.0,
    "needs_clarification": False,
    "clarifying_question": ""
}

def route_intent(llm: ChatOpenAI, user_query: str, artifact_store: ArtifactStore) -> dict:
    """
    사용자 질문의 의도를 분류하고 실행할 팩을 선택합니다.
    LLM은 숫자를 계산하지 않고, 오직 의도 분류와 팩 선택만 수행합니다.
    
    Returns:
        RouteResult 스키마에 맞는 dict
    """
    # 사용 가능한 팩 목록
    available_packs = [
        "EXPLAIN_RESULT", "COLUMN_DICTIONARY", "SINGLE_KPI", "STRUCTURE_SHARE",
        "ANOMALY_CHECK", "COMPARE_SIMPLE", "COMPARE_DEEP", "REPORT_PDF_SUMMARY",
        "EFFICIENCY_HINT", "EVIDENCE_LOOKUP"
    ]
    
    # 아티팩트 정보 수집
    tables = artifact_store.get_tables()
    table_info = []
    for artifact_id, df in tables[:3]:  # 최대 3개만
        table_info.append({
            "artifact_id": artifact_id,
            "rows": len(df),
            "columns": list(df.columns)[:10]  # 최대 10개 컬럼만
        })
    
    prompt = f"""
당신은 사용자 질문의 의도를 분류하고 적절한 분석 팩을 선택하는 라우터입니다.
숫자를 계산하거나 생성하지 마세요. 오직 의도 분류와 팩 선택만 수행하세요.

[사용 가능한 팩 목록]
{json.dumps(available_packs, ensure_ascii=False, indent=2)}

[사용자 질문]
{user_query}

[사용 가능한 데이터]
{json.dumps(table_info, ensure_ascii=False, indent=2) if table_info else "테이블 데이터 없음"}

[라우팅 스키마]
{json.dumps(ROUTE_RESULT_SCHEMA, ensure_ascii=False, indent=2)}

[규칙]
1. action은 EXPLAIN|ANALYZE|COMPARE|REPORT|OPTIMIZE 중 하나
2. domain은 logistics|inventory|settlement|sales|general 중 하나
3. pack은 위 팩 목록 중 하나
4. risk=high는 "정산/감사/증빙/청구/오차" 등 고위험 키워드 감지 시
5. confidence는 0.0~1.0 사이 값
6. 숫자 계산/생성 금지. params만 설정하세요.

[출력]
JSON만 출력하세요. 다른 텍스트는 출력하지 마세요.
"""
    
    try:
        messages = [
            SystemMessage(content="You are an intent router. Output JSON only, matching the RouteResult schema. Do not calculate numbers."),
            HumanMessage(content=prompt)
        ]
        resp = llm.invoke(messages).content
        result = _safe_json_load(resp)
        
        if not isinstance(result, dict):
            # 기본값 반환
            return {
                "action": "ANALYZE",
                "domain": "general",
                "pack": "EXPLAIN_RESULT",
                "params": {},
                "risk": "low",
                "confidence": 0.5,
                "needs_clarification": False,
                "clarifying_question": ""
            }
        
        # 기본값 보완
        result.setdefault("action", "ANALYZE")
        result.setdefault("domain", "general")
        result.setdefault("pack", "EXPLAIN_RESULT")
        result.setdefault("params", {})
        result.setdefault("risk", "low")
        result.setdefault("confidence", 0.5)
        result.setdefault("needs_clarification", False)
        result.setdefault("clarifying_question", "")
        
        return result
    except Exception as e:
        # 오류 시 기본값 반환
        return {
            "action": "ANALYZE",
            "domain": "general",
            "pack": "EXPLAIN_RESULT",
            "params": {},
            "risk": "low",
            "confidence": 0.3,
            "needs_clarification": True,
            "clarifying_question": f"의도 분류 중 오류 발생: {e}"
        }

def validate_analysis_plan(df: pd.DataFrame, plan: dict, user_query: str) -> dict:
    """
    LLM이 생성한 JSON 계획을 검증하고 수정합니다.
    LLM이 실수로 잘못된 필터나 top_n을 추가했을 경우 제거합니다.

    Args:
        df: 원본 DataFrame
        plan: LLM이 생성한 JSON 계획
        user_query: 사용자 질문

    Returns:
        수정된 JSON 계획
    """
    if not isinstance(plan, dict):
        return plan

    query_lower = (user_query or "").lower()
    modified = False

    # 1. 전체 합계 질문인데 필터가 있는 경우 → 필터 제거
    total_keywords = ["전체", "모든", "합계", "총액", "총합", "얼마"]
    exclude_keywords = ["만", "제외", "인 것", "보다"]

    is_total_query = any(k in query_lower for k in total_keywords)
    has_condition = any(k in query_lower for k in exclude_keywords)

    if is_total_query and not has_condition:
        if plan.get('filters') and len(plan['filters']) > 0:
            print(f"⚠️ [검증] 전체 합계 질문인데 필터가 있어서 제거: {plan['filters']}")
            plan['filters'] = []
            modified = True

    # 2. 상위 N개 요청이 없는데 top_n이 있는 경우 → top_n 제거
    top_keywords = ["상위", "top", "위까지", "많은", "적은"]
    has_top_request = any(k in query_lower for k in top_keywords)

    if not has_top_request:
        if 'top_n' in plan:
            print(f"⚠️ [검증] 상위 N개 요청이 없는데 top_n이 있어서 제거: top_n={plan['top_n']}")
            del plan['top_n']
            modified = True

    # 3. 그룹별 질문이 아닌데 groupby가 있는 경우 확인
    group_keywords = ["별", "각", "항목"]
    has_group_request = any(k in query_lower for k in group_keywords)

    if not has_group_request and plan.get('groupby') and len(plan['groupby']) > 0:
        print(f"⚠️ [검증] 그룹별 요청이 없는데 groupby가 있음: {plan['groupby']}")
        print(f"   → 단순 합계 질문으로 판단, groupby 제거")
        plan['groupby'] = []
        modified = True

    # 4. 결과 로깅
    if modified:
        print(f"✅ [검증] 계획 수정 완료")
    else:
        print(f"✅ [검증] 계획 이상 없음")

    return plan

# =========================================================
# Facts 구조 표준화 및 Claim Type 검증
# =========================================================
def build_standard_facts_dict(
    numeric_facts: dict = None,
    evidence_facts: List[dict] = None,
    perception_facts: List[dict] = None,
    meta: dict = None
) -> dict:
    """
    표준화된 facts dict 구조를 생성합니다.
    
    Args:
        numeric_facts: pandas 결정론적 결과
        evidence_facts: PDF/PPT/텍스트 근거 (출처 포함)
        perception_facts: OCR/이미지/차트 판독 (불확실)
        meta: 가정/제한사항
    
    Returns:
        표준화된 facts dict
    """
    return {
        "numeric_facts": numeric_facts or {},
        "evidence_facts": evidence_facts or [],
        "perception_facts": perception_facts or [],
        "meta": meta or {
            "assumptions": [],
            "limitations": []
        }
    }

def validate_facts_for_interpretation(facts: dict, claim_type: str, required_slots: List[str]) -> bool:
    """
    특정 Claim Type의 해석을 위해 필요한 facts 슬롯이 존재하는지 검증.
    
    Args:
        facts: 표준화된 facts dict
        claim_type: "Restate" | "Compare" | "ConditionalInsight"
        required_slots: 필요한 슬롯 이름 리스트 (예: ["total_today", "delta"])
    
    Returns:
        모든 슬롯이 존재하면 True
    """
    numeric = facts.get("numeric_facts", {})
    for slot in required_slots:
        # 중첩된 슬롯 지원 (예: "kpi_summary.total_today")
        if "." in slot:
            parts = slot.split(".")
            current = numeric
            for part in parts:
                if not isinstance(current, dict) or part not in current:
                    return False
                current = current[part]
        else:
            if slot not in numeric:
                return False
    return True

def build_interpretation_prompt(facts: dict, user_query: str, persona: str) -> str:
    """
    Facts-slot 기반 Claim Type A 제약이 포함된 LLM 해석 프롬프트 생성.
    
    Args:
        facts: 표준화된 facts dict
        user_query: 사용자 질문
        persona: 페르소나
    
    Returns:
        LLM 해석 프롬프트
    """
    numeric = facts.get("numeric_facts", {})
    evidence = facts.get("evidence_facts", [])
    perception = facts.get("perception_facts", [])
    meta = facts.get("meta", {})
    
    # 사용 가능한 facts 슬롯 목록 생성
    available_slots = []
    def _extract_slots(d: dict, prefix=""):
        for k, v in d.items():
            if isinstance(v, dict):
                _extract_slots(v, f"{prefix}{k}." if prefix else f"{k}.")
            elif isinstance(v, (int, float, str, bool)) or v is None:
                slot_name = f"{prefix}{k}" if prefix else k
                available_slots.append(slot_name)
    
    _extract_slots(numeric)
    
    prompt = f"""당신은 pandas가 계산한 facts를 기반으로 '해석 JSON'을 생성하는 AI입니다.

[절대 규칙 - 위반 금지]
1. 숫자 계산/생성 금지: 절대 숫자(0-9), 통화기호(₩,￦), %, 콤마를 출력하지 마세요.
2. Facts 슬롯 참조 필수: facts에 없는 슬롯을 요청(slot_keys)하면 안 됩니다.
3. Claim Type 제한: 아래 3가지 타입만 사용 가능합니다.
4. 출력은 반드시 JSON 1개만. JSON 외 텍스트 출력 금지.

[사용 가능한 Claim Type]
1) Restate (재진술)
   - numeric_facts의 값을 그대로 말함
   - 예: "오늘 총액은 {{total_today}}입니다."
   - 허용 표현: "~입니다", "~입니다", "~로 나타났습니다"
   - 금지: facts에 없는 값 추정/추론

2) Compare (비교)
   - numeric_facts 간 차이를 설명
   - 예: "어제 대비 {{delta}} 증가했습니다."
   - 허용 표현: "~대비 ~증가/감소", "~와 비교하여"
   - 금지: facts에 없는 비교 대상 사용

3) Conditional Insight (조건부 해석)
   - 반드시 대응되는 facts 슬롯이 존재해야 함
   - 예: "상위 품목 비중이 {{top_n_share}}%로 나타나, 관리 우선순위를 검토할 수 있습니다."
   - 허용 표현: "~로 나타나 ~할 수 있습니다", "~이므로 ~가능합니다"
   - 금지: facts 슬롯 없이 "집중되어 있다", "대부분", "급격히" 같은 정량적 암시

[사용 가능한 Facts 슬롯]
{json.dumps(available_slots, ensure_ascii=False, indent=2)}

[Perception Facts 사용 규칙]
- perception_facts는 설명 보조로만 사용
- "이미지에서 다음 문구가 인식됩니다 (OCR, 오류 가능)" 형식으로만 언급
- perception_facts를 근거로 수치/비율/비교 결론 생성 금지
- perception 정보를 numeric_facts와 결합해 해석하지 마세요

[사용자 질문]
{user_query}

[생성 규칙]
1. text_template에는 슬롯 참조만 포함하세요. 예: "오늘 총액은 {{kpi_summary.total_today}}입니다."
2. slot_keys에는 실제로 사용한 슬롯 키만 배열로 나열하세요.
3. 템플릿에 숫자/통화/%/콤마가 직접 포함되면 안 됩니다. (반드시 슬롯로만 표현)
4. "집중", "대부분", "급격히" 같은 정량적 암시 표현은 반드시 대응 facts 슬롯이 있을 때만 허용
5. risk=high 질문에서는 과장/단정 표현 금지

[출력 JSON 스키마]
{{\n  "claim_type": "Restate|Compare|ConditionalInsight",\n  "text_template": "...{{slot}}...",\n  "slot_keys": ["slot1","slot2"],\n  "provenance_refs": ["..."]\n}}

[출력 형식]
반드시 JSON 1개만 출력하세요. (JSON 외 텍스트 출력 금지)
"""
    
    if meta.get("limitations"):
        prompt += f"\n[제한사항]\n" + "\n".join(f"- {lim}" for lim in meta["limitations"])
    
    return prompt

# =========================================================
# 해석 JSON 스키마 및 렌더러 (슬롯 치환은 Python에서만)
# =========================================================
INTERPRETATION_JSON_SCHEMA = {
    "claim_type": "Restate|Compare|ConditionalInsight",
    "text_template": "...{{slot}}...",
    "slot_keys": ["slot1", "slot2"],
    "provenance_refs": ["..."]
}

# 팩별 허용 슬롯(명시적 화이트리스트).
# - PACK_ALLOWED_SLOT_KEYS: LLM이 요청 가능한 전체 slot_keys(스칼라+비스칼라 포함)
# - PACK_ALLOWED_TEMPLATE_SLOTS: 템플릿 치환(format)에 실제로 넣을 수 있는 스칼라 슬롯만
PACK_ALLOWED_SLOT_KEYS: Dict[str, set] = {
    "COMPARE_DEEP": {
        "kpi_summary.total_today",
        "kpi_summary.total_yesterday",
        "kpi_summary.delta",
        "kpi_summary.delta_rate",
        "top_n_share",
        "pareto_analysis.pareto_count",
        "pareto_analysis.pareto_share_percent",
        "pareto_analysis.pareto_items",
        # 아래는 dict/리스트(비스칼라): slot_keys로는 허용 가능하나 템플릿 치환은 금지
        "top_contributors_increase",
        "top_contributors_decrease",
        "share_by_item",
        "new_items",
        "discontinued_items",
        "numeric_conversion_failures",
    },
    "SINGLE_KPI": {
        "kpi.metric",
        "kpi.total",
        "kpi.count",
        "kpi.conversion_failures",
    },
    "EXPLAIN_RESULT": {
        "data_shape.rows",
        "data_shape.columns",
        "data_shape.column_names",
    },
}

PACK_ALLOWED_TEMPLATE_SLOTS: Dict[str, set] = {
    "COMPARE_DEEP": {
        "kpi_summary.total_today",
        "kpi_summary.total_yesterday",
        "kpi_summary.delta",
        "kpi_summary.delta_rate",
        "top_n_share",
        "pareto_analysis.pareto_count",
        "pareto_analysis.pareto_share_percent",
        # pareto_items는 list이므로 템플릿 치환 금지
    },
    "SINGLE_KPI": {
        "kpi.metric",
        "kpi.total",
        "kpi.count",
        "kpi.conversion_failures",
    },
    "EXPLAIN_RESULT": {
        "data_shape.rows",
        "data_shape.columns",
        # column_names는 list이므로 템플릿 치환 금지
    },
}

def _lookup_numeric_fact_slot(facts: dict, slot_key: str):
    """numeric_facts에서 slot_key(점 표기)를 따라 값을 가져옵니다. 없으면 None."""
    cur = facts.get("numeric_facts", {})
    for p in (slot_key or "").split("."):
        if not isinstance(cur, dict) or p not in cur:
            return None
        cur = cur[p]
    return cur

def _contains_banned_numeric_literals(s: str) -> bool:
    """
    템플릿에 '직접 수치'가 포함되면 True.
    - 숫자/통화/%는 금지
    - 콤마(,) 전체 금지는 과도하므로 허용
    - 단, 슬롯 플레이스홀더({{...}}) 내부의 숫자(예: top_1)는 검사에서 제외
    """
    if not isinstance(s, str):
        return True
    # 플레이스홀더 내부는 제거 후 검사 (slot key에 숫자가 있어도 OK)
    tmp = re.sub(r"\{\{[^{}]+\}\}", "", s)
    return bool(re.search(r"[0-9]|[₩￦%]", tmp))

def _extract_template_placeholders(text_template: str) -> List[str]:
    """{{slot}} 형태로 들어있는 slot key들을 추출합니다."""
    if not isinstance(text_template, str):
        return []
    keys = re.findall(r"\{\{([^{}]+)\}\}", text_template)
    return [k.strip() for k in keys if str(k).strip()]

def _sanitize_slot_placeholders(text_template: str) -> str:
    """
    text_template는 {{slot}} 형태를 사용하도록 강제.
    Python의 format은 {slot} 이므로, {{slot}} -> {slot}로 변환.
    """
    if not isinstance(text_template, str):
        return ""
    # {{slot}} -> {slot}
    return re.sub(r"\{\{([^{}]+)\}\}", r"{\1}", text_template)

def validate_interpretation_json(interpretation: dict, facts: dict, pack_name: str) -> dict:
    """
    LLM이 만든 해석 JSON을 검증/정규화합니다.
    - slot_keys는 PACK_ALLOWED_SLOT_KEYS에 없는 것은 제거
    - 템플릿에 사용된 플레이스홀더는 PACK_ALLOWED_TEMPLATE_SLOTS(스칼라)만 허용
    - ConditionalInsight는 (템플릿에 실제 사용된) 요구 슬롯이 모두 facts에 존재할 때만 유지, 아니면 Restate로 강등
    - text_template에 숫자 리터럴이 있으면 렌더링 차단(근거부족으로 강등)
    """
    if not isinstance(interpretation, dict):
        return {
            "claim_type": "Restate",
            "text_template": "근거 부족: 해석 JSON이 유효하지 않습니다.",
            "slot_keys": [],
            "provenance_refs": []
        }

    claim_type = interpretation.get("claim_type") or "Restate"
    text_template = interpretation.get("text_template") or ""
    slot_keys = interpretation.get("slot_keys") or []
    prov_refs = interpretation.get("provenance_refs") or []

    if claim_type not in {"Restate", "Compare", "ConditionalInsight"}:
        claim_type = "Restate"

    if not isinstance(slot_keys, list):
        slot_keys = []
    slot_keys = [str(k) for k in slot_keys if isinstance(k, (str, int, float))]

    allowed_keys = PACK_ALLOWED_SLOT_KEYS.get(pack_name, set())
    allowed_template = PACK_ALLOWED_TEMPLATE_SLOTS.get(pack_name, set())
    slot_keys = [k for k in slot_keys if k in allowed_keys]

    # 템플릿 직접 수치 차단 (LLM이 숫자를 써도 최종 렌더링에 반영되지 않도록)
    if _contains_banned_numeric_literals(text_template):
        return {
            "claim_type": "Restate",
            "text_template": "근거 부족: 템플릿에 직접 수치가 포함되어 렌더링하지 않습니다.",
            "slot_keys": [],
            "provenance_refs": prov_refs if isinstance(prov_refs, list) else []
        }

    # 템플릿 플레이스홀더 검증: 템플릿 치환은 스칼라 슬롯만 허용
    placeholders = _extract_template_placeholders(text_template)
    if placeholders:
        # 템플릿에 등장하는 slot은 slot_keys에 포함되어야 함
        missing_decl = [p for p in placeholders if p not in slot_keys]
        if missing_decl:
            claim_type = "Restate"
            text_template = "근거 부족: 템플릿이 참조하는 slot이 slot_keys에 선언되지 않았습니다."
            slot_keys = []
        else:
            # 템플릿에 등장하는 slot은 TEMPLATE_SLOTS(스칼라)만 허용
            not_scalar = [p for p in placeholders if p not in allowed_template]
            if not_scalar:
                claim_type = "Restate"
                text_template = "근거 부족: 템플릿 치환이 허용되지 않는(비스칼라) 슬롯을 참조합니다."
                slot_keys = []
            else:
                # slot_keys는 템플릿에 실제 사용된 slot만 남김(불필요한 삽입 위험 감소)
                slot_keys = placeholders

    # ConditionalInsight는 요구 슬롯이 모두 있을 때만
    if claim_type == "ConditionalInsight":
        # 요구 슬롯 = 템플릿에 실제 사용된 slot_keys
        missing = [k for k in slot_keys if _lookup_numeric_fact_slot(facts, k) is None]
        if missing:
            claim_type = "Restate"
            text_template = "근거 부족: 조건부 해석에 필요한 facts 슬롯이 부족합니다."
            slot_keys = []

    return {
        "claim_type": claim_type,
        "text_template": text_template,
        "slot_keys": slot_keys,
        "provenance_refs": prov_refs if isinstance(prov_refs, list) else []
    }

def generate_interpretation(llm: ChatOpenAI, facts: dict, user_query: str, persona: str, pack_name: str) -> dict:
    """
    Facts-slot 기반 Claim Type A 제약으로 LLM 해석 JSON 생성.
    - LLM은 slot_keys만 반환 (실제 값/숫자 반환 금지)
    - JSON 외 텍스트 출력 금지
    
    Args:
        llm: ChatOpenAI 인스턴스
        facts: 표준화된 facts dict
        user_query: 사용자 질문
        persona: 페르소나
        pack_name: 팩 이름 (ALLOWED_SLOTS 검증용)
    
    Returns:
        해석 JSON dict (slot_keys 기반)
    """
    prompt = build_interpretation_prompt(facts, user_query, persona)
    allowed = sorted(list(PACK_ALLOWED_SLOT_KEYS.get(pack_name, set())))
    
    messages = [
        SystemMessage(content="You must output JSON only. Never output numeric literals. Return slot_keys only; do not return slot values."),
        HumanMessage(content=prompt + "\n\n[ALLOWED_SLOTS]\n" + json.dumps(allowed, ensure_ascii=False, indent=2))
    ]
    
    try:
        response = llm.invoke(messages).content
        parsed = _safe_json_load(response)
        validated = validate_interpretation_json(parsed, facts, pack_name)
        return validated
    except Exception as e:
        return {
            "claim_type": "Restate",
            "text_template": f"근거 부족: 해석 생성 중 오류 발생 ({e})",
            "slot_keys": [],
            "provenance_refs": []
        }

def render_interpretation(interpretation: dict, facts: dict, pack_name: str) -> str:
    """
    해석 JSON을 실제 텍스트로 렌더링합니다.
    - 슬롯 치환은 Python에서만 수행
    - slot_keys에 없는 값은 절대 렌더링 금지
    - slot_keys가 facts에 없으면 '근거 부족'으로 강등/오류 표기
    """
    it = validate_interpretation_json(interpretation, facts, pack_name)
    text_template = _sanitize_slot_placeholders(it.get("text_template", ""))
    slot_keys = it.get("slot_keys", [])

    # 실제 치환은 TEMPLATE_SLOTS(스칼라)만 허용
    allowed_template = PACK_ALLOWED_TEMPLATE_SLOTS.get(pack_name, set())
    slot_keys_for_template = [k for k in slot_keys if k in allowed_template]

    # slot_keys -> facts lookup (템플릿 치환 대상만)
    slots = {k: _lookup_numeric_fact_slot(facts, k) for k in slot_keys_for_template}
    missing = [k for k, v in slots.items() if v is None]
    if missing:
        return "근거 부족: 요청된 facts 슬롯을 찾을 수 없습니다. (" + ", ".join(missing) + ")"

    try:
        rendered = text_template.format(**slots) if text_template else "근거 부족: 템플릿이 비어 있습니다."
    except Exception as e:
        rendered = f"근거 부족: 템플릿 렌더링 오류 ({e})"

    # provenance refs는 텍스트로만 덧붙임 (수치 근거로 재해석 금지)
    prov_refs = it.get("provenance_refs") or []
    if prov_refs:
        rendered += "\n\n[근거]\n- " + "\n- ".join([str(x) for x in prov_refs[:8]])
    return rendered

# =========================================================
# Pack Dispatch: 팩 실행 시스템
# =========================================================
def dispatch_pack(route_result: dict, artifact_store: ArtifactStore) -> Tuple[dict, str]:
    """
    RouteResult에 따라 적절한 팩을 실행합니다.
    모든 숫자 계산은 pandas로만 수행하며, 표준화된 facts dict를 반환합니다.
    
    Returns:
        (standardized_facts_dict, provenance_string)
        standardized_facts_dict는 build_standard_facts_dict() 구조를 따름
    """
    pack_name = route_result.get("pack", "EXPLAIN_RESULT")
    params = route_result.get("params", {})
    
    # 테이블 데이터 가져오기
    tables = artifact_store.get_tables()
    if not tables:
        return build_standard_facts_dict(
            meta={"limitations": ["테이블 데이터가 없습니다."]}
        ), ""
    
    # 첫 번째 테이블 사용 (향후 다중 테이블 지원 확장 가능)
    artifact_id, df = tables[0]
    provenance = f"[근거: {artifact_id}]"
    
    if pack_name == "COMPARE_DEEP":
        return pack_compare_deep(df, params, artifact_store)
    elif pack_name == "SINGLE_KPI":
        return pack_single_kpi(df, params)
    elif pack_name == "EXPLAIN_RESULT":
        return pack_explain_result(df, params)
    elif pack_name == "EVIDENCE_LOOKUP":
        return pack_evidence_lookup(artifact_store, params)
    else:
        # 기본: 단순 설명
        return pack_explain_result(df, params)

def pack_compare_deep(df: pd.DataFrame, params: dict, artifact_store: ArtifactStore) -> Tuple[dict, str]:
    """
    COMPARE_DEEP v1.5: 어제 vs 오늘 심층 비교 팩
    
    [생성하는 Facts 슬롯]
    numeric_facts:
      - kpi_summary.total_today: float (오늘 총액)
      - kpi_summary.total_yesterday: float | None (어제 총액, v1.5는 단일 파일이므로 None 가능)
      - kpi_summary.delta: float | None (차이)
      - kpi_summary.delta_rate: float | None (변화율 %)
      - top_contributors_increase: dict (증가 기여도 TOP N)
      - top_contributors_decrease: dict (감소 기여도 TOP N)
      - pareto_analysis.pareto_count: int (Pareto 80% 기여 품목 수)
      - pareto_analysis.pareto_share_percent: float (Pareto 80% 기여 비중)
      - pareto_analysis.pareto_items: List[str] (Pareto 품목 목록)
      - share_by_item: dict (품목별 비중 %)
      - top_n_share: float (상위 N개 비중 %)
      - new_items: List[str] (신규 품목)
      - discontinued_items: List[str] (소멸 품목)
      - numeric_conversion_failures: dict (숫자 변환 실패 정보)
    
    [해석 허용 범위]
    - Restate: 모든 numeric_facts 값 재진술 가능
    - Compare: total_yesterday와 total_today가 모두 있으면 delta/delta_rate 비교 가능
    - Conditional Insight:
      * pareto_share_percent가 있으면 "집중도" 관련 표현 가능
      * top_n_share가 있으면 "상위 N개 비중" 표현 가능
      * delta_rate가 있으면 "증가/감소율" 표현 가능
      * facts 슬롯 없으면 "집중", "대부분", "급격히" 같은 표현 금지
    
    모든 계산은 pandas 결정론적.
    """
    provenance_parts = []
    
    # 컬럼 매핑 (best-effort)
    col_item = _pick_col(df, ["상품명", "품목명", "item", "name", "품목"])
    col_date = _pick_col(df, ["날짜", "일자", "date", "입고일자", "기준일자"])
    col_amount = _pick_col(df, ["금액", "총액", "amount", "매출", "정산금액", "불용재고 총액", "정상재고 총액"])
    col_qty = _pick_col(df, ["수량", "qty", "quantity", "전산수량"])
    
    if not col_item:
        return build_standard_facts_dict(
            meta={"limitations": ["품목명 컬럼을 찾을 수 없습니다."]}
        ), ""
    
    # 숫자 컬럼 정규화 및 변환 실패 추적
    numeric_failures = {}
    work = df.copy()
    
    for col in [col_amount, col_qty]:
        if col and col in work.columns:
            original = work[col].copy()
            coerced = _coerce_numeric_series(work[col])
            work[col] = coerced
            
            # 변환 실패 추적
            failed_mask = pd.isna(coerced) & original.notna()
            if failed_mask.any():
                failed_count = int(failed_mask.sum())
                failed_samples = original[failed_mask].head(5).tolist()
                numeric_failures[col] = {
                    "count": failed_count,
                    "samples": [str(v) for v in failed_samples]
                }
    
    if numeric_failures:
        provenance_parts.append(f"숫자 변환 실패: {len(numeric_failures)}개 컬럼")
    
    compare_target = params.get("compare_target", "yesterday_vs_today")
    top_n = params.get("top_n", 10)
    
    # 메인 메트릭 선택 (금액 우선, 없으면 수량)
    main_metric_col = col_amount if col_amount and col_amount in work.columns else col_qty
    if not main_metric_col or main_metric_col not in work.columns:
        return build_standard_facts_dict(
            meta={"limitations": ["비교할 숫자 컬럼을 찾을 수 없습니다."]}
        ), ""
    
    # KPI 요약: 전체 합계
    total_today = float(_coerce_numeric_series(work[main_metric_col]).sum(skipna=True))
    
    # numeric_facts 구성
    numeric_facts = {
        "kpi_summary": {
            "total_today": float(total_today),
            "total_yesterday": None,  # v1.5에서는 단일 파일 기준
            "delta": None,
            "delta_rate": None
        }
    }
    
    if numeric_failures:
        numeric_facts["numeric_conversion_failures"] = numeric_failures
    
    # 그룹별 집계 (품목별)
    if col_item in work.columns:
        grouped = work.groupby(col_item, dropna=False).agg({
            main_metric_col: 'sum'
        }).reset_index()
        grouped = grouped.sort_values(main_metric_col, ascending=False)
        
        # 증감 기여도 TOP N (증가)
        top_increase = grouped.head(top_n).copy()
        top_increase_dict = {
            f"top_{i+1}": {
                "item": str(k),
                "value": float(v)
            }
            for i, (k, v) in enumerate(top_increase.set_index(col_item)[main_metric_col].to_dict().items())
        }
        
        # 증감 기여도 TOP N (감소)
        top_decrease = grouped.tail(top_n).copy()
        top_decrease_dict = {
            f"bottom_{i+1}": {
                "item": str(k),
                "value": float(v)
            }
            for i, (k, v) in enumerate(top_decrease.set_index(col_item)[main_metric_col].to_dict().items())
        }
        
        # Pareto 분석 (집중도) - index 혼동 없는 방식
        # 1) metric 기준 정렬한 DataFrame g 생성 후 reset_index(drop=True)
        g = grouped[[col_item, main_metric_col]].copy()
        g = g.sort_values(main_metric_col, ascending=False).reset_index(drop=True)
        vals = g[main_metric_col].fillna(0)
        cumsum = vals.cumsum()
        total_sum = float(vals.sum())
        pareto_threshold = total_sum * 0.8

        if total_sum > 0 and len(g) > 0:
            mask = cumsum <= pareto_threshold
            # Pareto는 최소 1개 품목은 포함되도록 보정
            if mask.sum() == 0:
                mask.iloc[0] = True
            pareto_items = g.loc[mask, col_item].astype(str).tolist()
            pareto_count = int(len(pareto_items))
            pareto_sum = float(vals.loc[mask].sum())
            pareto_share = float((pareto_sum / total_sum * 100) if total_sum > 0 else 0)
        else:
            pareto_items = []
            pareto_count = 0
            pareto_share = 0.0
        
        # 비중 계산
        grouped['share'] = (grouped[main_metric_col] / total_today * 100) if total_today > 0 else 0
        share_dict = {str(k): float(v) for k, v in grouped.set_index(col_item)['share'].to_dict().items()}
        
        # 상위 N개 비중 계산 (해석 슬롯)
        top_n_share = float(grouped.head(top_n)[main_metric_col].sum() / total_today * 100) if total_today > 0 else 0.0
        
        numeric_facts.update({
            "top_contributors_increase": top_increase_dict,
            "top_contributors_decrease": top_decrease_dict,
            "pareto_analysis": {
                "pareto_count": pareto_count,
                "pareto_share_percent": float(pareto_share),
                "pareto_items": pareto_items[:10]
            },
            "share_by_item": dict(list(share_dict.items())[:20]),  # 최대 20개만
            "top_n_share": top_n_share,  # 해석 슬롯: 상위 N개 비중
            "new_items": [],  # v1.5에서는 단일 파일이므로 빈 리스트
            "discontinued_items": []
        })
    
    # evidence_facts: 관련 문서 검색 (선택적)
    evidence_facts = []
    query_hint = params.get("key_hint", [])
    if query_hint:
        search_query = " ".join(query_hint[:3])
        evidence_artifacts = artifact_store.search_text(search_query, max_results=3)
        evidence_facts = [
            {
                "source": a.file_name,
                "location": a.source_loc,
                "text_snippet": a.text[:300],
                "provenance": a.provenance
            }
            for a in evidence_artifacts
        ]
    
    # perception_facts: 이미지/OCR (현재는 빈 리스트, 향후 확장)
    perception_facts = []
    
    # meta: 가정 및 제한사항
    meta = {
        "assumptions": [
            "v1.5는 단일 파일 기준 비교입니다",
            "날짜 기반 비교는 향후 두 파일/날짜 컬럼 필요"
        ],
        "limitations": []
    }
    if numeric_failures:
        meta["limitations"].append(f"{len(numeric_failures)}개 컬럼에서 숫자 변환 실패")
    
    facts = build_standard_facts_dict(
        numeric_facts=numeric_facts,
        evidence_facts=evidence_facts,
        perception_facts=perception_facts,
        meta=meta
    )
    
    provenance = "; ".join(provenance_parts) if provenance_parts else f"[근거: {main_metric_col} 집계]"
    
    return facts, provenance

def pack_single_kpi(df: pd.DataFrame, params: dict) -> Tuple[dict, str]:
    """
    단일 KPI 계산 팩.
    
    [생성하는 Facts 슬롯]
    numeric_facts:
      - kpi.metric: str (메트릭 컬럼명)
      - kpi.total: float (총합)
      - kpi.count: int (행 수)
      - kpi.conversion_failures: int (변환 실패 개수)
    
    [해석 허용 범위]
    - Restate: kpi.total, kpi.count 재진술 가능
    - Compare: 없음 (단일 값)
    - Conditional Insight: 없음
    """
    metric_hint = params.get("metric_hint", [])
    if not metric_hint:
        metric_hint = ["amount", "금액", "총액"]
    
    # 메트릭 컬럼 찾기
    metric_col = None
    for hint in metric_hint:
        metric_col = _pick_col(df, [hint])
        if metric_col:
            break
    
    if not metric_col:
        # 첫 번째 숫자 컬럼 사용
        numeric_cols = _infer_numeric_columns(df)
        if numeric_cols:
            metric_col = numeric_cols[0]
        else:
            return build_standard_facts_dict(
                meta={"limitations": ["KPI 계산할 컬럼을 찾을 수 없습니다."]}
            ), ""
    
    # 숫자 변환 및 실패 추적
    coerced = _coerce_numeric_series(df[metric_col])
    total = float(coerced.sum(skipna=True))
    
    failed_count = int(pd.isna(coerced).sum())
    
    numeric_facts = {
        "kpi": {
            "metric": metric_col,
            "total": total,
            "count": len(df),
            "conversion_failures": failed_count
        }
    }
    
    provenance = f"[근거: {metric_col} 집계]"
    if failed_count > 0:
        provenance += f" (변환 실패: {failed_count}개)"
    
    return build_standard_facts_dict(
        numeric_facts=numeric_facts,
        meta={"limitations": [f"변환 실패: {failed_count}개"] if failed_count > 0 else []}
    ), provenance

def pack_explain_result(df: pd.DataFrame, params: dict) -> Tuple[dict, str]:
    """
    결과 설명 팩 (기본).
    
    [생성하는 Facts 슬롯]
    numeric_facts:
      - data_shape.rows: int (행 수)
      - data_shape.columns: int (컬럼 수)
      - data_shape.column_names: List[str] (컬럼명 목록)
    
    [해석 허용 범위]
    - Restate: data_shape 값 재진술 가능
    - Compare: 없음
    - Conditional Insight: 없음
    """
    numeric_facts = {
        "data_shape": {
            "rows": len(df),
            "columns": len(df.columns),
            "column_names": list(df.columns)
        }
    }
    return build_standard_facts_dict(numeric_facts=numeric_facts), "[근거: 데이터 구조]"

def pack_evidence_lookup(artifact_store: ArtifactStore, params: dict) -> Tuple[dict, str]:
    """
    문서 근거 검색 팩.
    
    [생성하는 Facts 슬롯]
    evidence_facts:
      - evidence_results: List[dict] (검색된 문서 근거 목록)
        각 항목: artifact_id, file_name, source_loc, provenance, text_snippet
    
    [해석 허용 범위]
    - Restate: evidence_results 개수/출처 재진술 가능
    - Compare: 없음
    - Conditional Insight: 없음 (evidence는 설명 보조용)
    """
    query = params.get("query", "")
    if not query:
        return build_standard_facts_dict(
            meta={"limitations": ["검색어가 필요합니다."]}
        ), ""
    
    results = artifact_store.search_text(query, max_results=5)
    evidence_facts = [
        {
            "artifact_id": a.artifact_id,
            "file_name": a.file_name,
            "source_loc": a.source_loc,
            "provenance": a.provenance,
            "text_snippet": a.text[:500]  # 최대 500자
        }
        for a in results
    ]
    
    return build_standard_facts_dict(evidence_facts=evidence_facts), f"[근거: {len(results)}개 문서에서 검색]"

def execute_analysis_plan(df: pd.DataFrame, plan: dict) -> Tuple[pd.DataFrame, Optional[plt.Figure], str, dict]:
    """
    Returns: (result_df, fig_or_none, provenance_text, conversion_failures_dict)
    provenance_text should be used as "근거" in UI.
    conversion_failures_dict는 숫자 변환 실패 정보를 담은 dict (조용히 무시 금지).
    """
    work = df.copy()
    provenance = f"[근거: {plan.get('dataset_file','(unknown)')}]"
    conversion_failures = {}  # 변환 실패 추적

    # optional select columns (but keep needed cols)
    sel_cols = plan.get("select_columns") or []
    # apply filters first (needs columns)
    work = _apply_filters(work, plan.get("filters", []))

    groupby_cols = plan.get("groupby") or []
    metrics = plan.get("metrics") or []

    # build agg dict 및 숫자 변환 (통일된 _coerce_numeric_series 사용)
    agg_map = {}
    metric_names = []
    for m in metrics:
        col = m.get("column")
        agg = (m.get("agg") or "").lower()
        if col in work.columns and agg in {"sum","mean","count","min","max"}:
            # 숫자 변환: _coerce_numeric_series 사용 (통일)
            original = work[col].copy()
            coerced = _coerce_numeric_series(work[col])
            
            # 변환 실패 추적
            failed_mask = pd.isna(coerced) & original.notna()
            if failed_mask.any():
                failed_count = int(failed_mask.sum())
                failed_samples = original[failed_mask].head(5).tolist()
                conversion_failures[col] = {
                    "count": failed_count,
                    "samples": [str(v) for v in failed_samples],
                    "agg": agg
                }
            
            # 변환된 값으로 교체
            work[col] = coerced
            agg_map[col] = agg
            metric_names.append(f"{agg}({col})")

    if groupby_cols:
        gb_cols = [c for c in groupby_cols if c in work.columns]
        if not gb_cols:
            # fallback: no valid groupby
            res = work.copy()
        else:
            if agg_map:
                res = work.groupby(gb_cols, dropna=False).agg(agg_map).reset_index()
                # rename metric columns to agg(col)
                rename = {col: f"{agg_map[col]}({col})" for col in agg_map}
                res = res.rename(columns=rename)
            else:
                # default metric: count
                res = work.groupby(gb_cols, dropna=False).size().reset_index(name="count")
                metric_names = ["count"]
    else:
        # no groupby: overall aggregation
        if agg_map:
            row = {}
            for col, agg in agg_map.items():
                # 이미 _coerce_numeric_series로 변환된 상태
                s = work[col]
                if agg == "sum":
                    row[f"sum({col})"] = float(s.sum(skipna=True))
                elif agg == "mean":
                    row[f"mean({col})"] = float(s.mean(skipna=True))
                elif agg == "count":
                    row[f"count({col})"] = int(s.count())
                elif agg == "min":
                    row[f"min({col})"] = float(s.min(skipna=True))
                elif agg == "max":
                    row[f"max({col})"] = float(s.max(skipna=True))
            res = pd.DataFrame([row])
            metric_names = list(row.keys())
        else:
            res = pd.DataFrame({"count": [len(work)]})
            metric_names = ["count"]
    
    # 변환 실패 정보를 provenance에 추가
    if conversion_failures:
        failure_summary = "; ".join([f"{col}({info['count']}개 실패)" for col, info in conversion_failures.items()])
        provenance += f" | 변환 실패: {failure_summary}"

    # sorting
    sort = plan.get("sort") or {}
    by = sort.get("by")
    asc = bool(sort.get("ascending", False))
    if by and by in res.columns:
        res = res.sort_values(by=by, ascending=asc)

    # top_n
    top_n = plan.get("top_n")
    if isinstance(top_n, int) and top_n > 0 and len(res) > top_n:
        res = res.head(top_n)

    # chart
    fig = None
    chart = plan.get("chart") or {}
    ctype = (chart.get("type") or "").lower()
    x = chart.get("x")
    y = chart.get("y")
    title = chart.get("title") or ""
    xlabel = chart.get("xlabel") or (x or "")
    ylabel = chart.get("ylabel") or (y or "")

    if ctype in {"bar","line","pie"} and y and y in res.columns:
        fig, ax = plt.subplots()
        if ctype == "pie":
            # pie uses y values and x labels if provided
            labels = res[x].astype(str).tolist() if x and x in res.columns else res.index.astype(str).tolist()
            ax.pie(res[y].fillna(0).tolist(), labels=labels)
        else:
            xv = res[x] if x and x in res.columns else res.index
            if ctype == "bar":
                ax.bar(xv.astype(str), res[y].fillna(0))
                ax.tick_params(axis='x', labelrotation=45)
            elif ctype == "line":
                ax.plot(xv, res[y].fillna(0))
            ax.set_xlabel(xlabel)
            ax.set_ylabel(ylabel)
        if title:
            ax.set_title(title)
        fig.tight_layout()

    return res, fig, provenance

# =========================================================
# Streamlit app
# =========================================================
def init_state():
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "file_cache" not in st.session_state:
        st.session_state.file_cache = {}
    if "processed_file_names" not in st.session_state:
        st.session_state.processed_file_names = set()
    if "artifact_store" not in st.session_state:
        st.session_state.artifact_store = ArtifactStore()

def sidebar():
    st.sidebar.title("⚙️ 설정")

    api_key = st.sidebar.text_input(
        "OpenAI API Key",
        type="password",
        value=os.getenv("OPENAI_API_KEY", ""),
        help="키는 코드/저장소에 저장되지 않습니다. 현재 세션에서만 사용합니다."
    )
    temperature = st.sidebar.slider("Temperature", 0.0, 1.0, 0.2, 0.1)
    persona = st.sidebar.selectbox("페르소나", ["물류 비서", "감사관(정산 검증)", "기획자(PPT 요약)"], index=0)

    st.sidebar.divider()
    mode = st.sidebar.radio("모드", ["Q&A(검색/요약)", "분석(표/그래프)"], index=0)
    use_web = st.sidebar.checkbox("외부 검색(웹) 사용", value=False, help="기업 데모에서는 기본 OFF 권장.")

    st.sidebar.divider()

    # File upload section in sidebar
    st.sidebar.subheader("📎 파일")
    uploaded_files = st.sidebar.file_uploader(
        "파일 업로드",
        type=["xlsx", "csv", "pdf", "pptx", "docx", "png", "jpg", "jpeg", "webp"],
        accept_multiple_files=True,
        label_visibility="collapsed"
    )

    if uploaded_files:
        for uf in uploaded_files:
            if uf.name in st.session_state.processed_file_names:
                continue
            with st.spinner(f"📥 분석 중: {uf.name}"):
                cached = parse_uploaded_file(uf)
                st.session_state.file_cache[uf.name] = cached
                st.session_state.processed_file_names.add(uf.name)
                # ArtifactStore에도 추가
                artifact_ids = st.session_state.artifact_store.add_from_cached_file(cached)

    # Show uploaded files compactly
    if st.session_state.file_cache:
        st.sidebar.caption(f"업로드된 파일 ({len(st.session_state.file_cache)}개)")
        for fn, cf in st.session_state.file_cache.items():
            if cf.kind == "table":
                r, c = cf.meta.get("shape", [0, 0])
                st.sidebar.caption(f"📊 {fn[:20]}... ({r}x{c})")
            elif cf.kind == "pdf":
                pages = cf.meta.get('pages', 0)
                st.sidebar.caption(f"📄 {fn[:20]}... ({pages}p)")
            elif cf.kind == "ppt":
                slides = cf.meta.get('slides', 0)
                st.sidebar.caption(f"📊 {fn[:20]}... ({slides}s)")
            elif cf.kind == "docx":
                st.sidebar.caption(f"📝 {fn[:20]}...")
            elif cf.kind == "image":
                st.sidebar.caption(f"🖼️ {fn[:20]}...")

    return api_key, temperature, persona, mode, use_web

def file_uploader_area():
    st.subheader("📎 파일 업로드 (엑셀/PDF/PPT/DOCX/사진)")
    uploaded_files = st.file_uploader(
        "여러 파일을 동시에 업로드할 수 있습니다.",
        type=["xlsx", "csv", "pdf", "pptx", "docx", "png", "jpg", "jpeg", "webp"],
        accept_multiple_files=True,
    )

    if uploaded_files:
        for uf in uploaded_files:
            if uf.name in st.session_state.processed_file_names:
                continue
            with st.spinner(f"📥 분석 중: {uf.name}"):
                cached = parse_uploaded_file(uf)
                st.session_state.file_cache[uf.name] = cached
                st.session_state.processed_file_names.add(uf.name)
                # ArtifactStore에도 추가
                artifact_ids = st.session_state.artifact_store.add_from_cached_file(cached)

    if st.session_state.file_cache:
        st.caption("✅ 분석 완료된 파일")
        for fn, cf in st.session_state.file_cache.items():
            if cf.kind == "table":
                r, c = cf.meta.get("shape",[0,0])
                st.write(f"- {fn}  (TABLE, {r}x{c})")
            elif cf.kind == "pdf":
                st.write(f"- {fn}  (PDF, pages={cf.meta.get('pages',0)})")
            elif cf.kind == "ppt":
                st.write(f"- {fn}  (PPT, slides={cf.meta.get('slides',0)})")
            elif cf.kind == "docx":
                st.write(f"- {fn}  (DOCX)")
            elif cf.kind == "image":
                st.write(f"- {fn}  (IMAGE)")
            else:
                st.write(f"- {fn}  (UNKNOWN)")
    else:
        st.info("업로드된 파일이 아직 없습니다.")

def chat_history_render():
    for m in st.session_state.messages:
        with st.chat_message(m["role"]):
            st.markdown(m["content"])

def run_qa(llm: ChatOpenAI, persona: str, use_web: bool, selected_files: List[CachedFile]):
    # Render history
    if len(st.session_state.messages) == 0:
        # Welcome message when no chat history
        st.markdown("""
        <div style="text-align: center; padding: 3rem 1rem; color: #666;">
            <h2 style="margin-bottom: 1rem;">💬 무엇을 도와드릴까요?</h2>
            <p style="margin-bottom: 2rem;">파일을 업로드하고 질문해보세요</p>
        </div>
        """, unsafe_allow_html=True)
    else:
        chat_history_render()

    prompt = st.chat_input("메시지를 입력하세요...")
    if not prompt:
        return

    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    with st.chat_message("assistant"):
        placeholder = st.empty()
        try:
            system_prompt = build_system_prompt(persona, mode="qa")
            # ---- Pandas-first numeric analysis guardrail ----
            # If the user asks for totals/aggregations, do NOT let LLM do arithmetic.
            # Instead: create an analysis plan JSON -> execute with pandas -> show deterministic results.
            if _is_numeric_query(prompt) or _is_report_query(prompt):
                table_files = [f for f in selected_files if f.kind == "table" and f.df is not None and not f.df.empty]
                if table_files:
                    # choose first table by default (PoV stability). If multiple, user should select in analysis mode.
                    dataset_file = table_files[0].name
                    df0 = table_files[0].df
                    # ---- Report shortcut (no LLM arithmetic) ----
                    # Report/summary requests must be deterministic.
                    if _is_report_query(prompt):
                        md = build_inventory_report_markdown(df0, dataset_file)
                        st.markdown(md)
                        st.session_state.messages.append({"role": "assistant", "content": md})
                        return
                    # ---- End report shortcut ----
                    with st.spinner("pandas로 숫자 집계 중..."):
                        try:
                            plan_prompt = _analysis_prompt(df0, dataset_file, prompt)
                            plan_msgs = [
                                SystemMessage(content=build_system_prompt(persona, mode="analysis")),
                                HumanMessage(content=plan_prompt),
                            ]
                            plan_resp = llm.invoke(plan_msgs).content
                            plan = _safe_json_load(plan_resp)
                            if isinstance(plan, dict):
                                # Validate and fix LLM-generated plan (prevent wrong filters/top_n)
                                plan = validate_analysis_plan(df0, plan, prompt)
                                plan["dataset_file"] = dataset_file
                                result_df, fig, provenance, conversion_failures = execute_analysis_plan(df0, plan)

                                # 변환 실패 정보 표시 (조용히 무시 금지)
                                if conversion_failures:
                                    with st.expander("⚠️ 숫자 변환 실패 정보", expanded=False):
                                        for col, info in conversion_failures.items():
                                            st.write(f"**{col}**: {info['count']}개 실패")
                                            st.caption(f"샘플: {', '.join(info['samples'][:3])}")

                                # Determine if visualization is needed
                                needs_viz = _needs_visualization(prompt)

                                # Always show text summary
                                st.markdown(_format_result_markdown(result_df, provenance))

                                # Only render table/chart if visualization is actually needed
                                if needs_viz:
                                    if len(result_df) > 1 or len(result_df.columns) > 8:
                                        st.dataframe(result_df, use_container_width=True)
                                    if fig is not None:
                                        st.pyplot(fig, clear_figure=True)

                                # Save to history (text only; numbers are from pandas)
                                resp_text = _format_result_markdown(result_df, provenance)
                                st.session_state.messages.append({"role": "assistant", "content": resp_text})
                                return
                        except Exception as e:
                            st.warning(f"숫자 집계 경로에서 오류가 발생해 일반 Q&A로 fallback 합니다: {e}")
                # if no table files, continue with normal QA
            # ---- End guardrail ----

            context_text, context_images = build_context_blocks(selected_files, prompt)

            final_text = "=== 참고 자료(선택된 파일에서 추출) ===\n" + context_text + "\n\n=== 사용자 질문 ===\n" + prompt
            messages = [SystemMessage(content=system_prompt)]

            history = st.session_state.messages[-10:-1]
            for h in history:
                if h["role"] == "user":
                    messages.append(HumanMessage(content=h["content"]))
                else:
                    messages.append(AIMessage(content=h["content"]))

            human_content = [{"type": "text", "text": final_text}]
            for img in context_images:
                human_content.append(img)
            messages.append(HumanMessage(content=human_content))

            # Streaming response for ChatGPT-like experience
            full_response = ""
            for chunk in llm.stream(messages):
                if hasattr(chunk, 'content') and chunk.content:
                    full_response += chunk.content
                    placeholder.markdown(full_response + "▌")  # Cursor effect

            # Final update without cursor
            placeholder.markdown(full_response)
            st.session_state.messages.append({"role": "assistant", "content": full_response})
        except Exception as e:
            st.error(f"오류: {e}")


# =========================================================
# Simple PDF report export (1-page) for PoV
# =========================================================
def _fig_to_png_bytes(fig: Optional[plt.Figure]) -> Optional[bytes]:
    if fig is None:
        return None
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    return buf.getvalue()

def build_pdf_report_bytes(
    title: str,
    subtitle: str,
    provenance: str,
    result_df: pd.DataFrame,
    fig_png: Optional[bytes],
    extra_notes: str = ""
) -> bytes:
    """
    Build a one-page PDF containing: title, subtitle, provenance, small table, and optional chart image.
    """
    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    width, height = A4

    # Margins
    x0 = 1.5 * cm
    y = height - 1.5 * cm

    # Title
    c.setFont("Helvetica-Bold", 14)
    c.drawString(x0, y, _truncate(title, 120))
    y -= 0.7 * cm

    c.setFont("Helvetica", 10)
    c.drawString(x0, y, _truncate(subtitle, 160))
    y -= 0.6 * cm

    # Provenance
    c.setFont("Helvetica", 9)
    prov = _truncate(provenance or "", 260)
    if prov:
        c.drawString(x0, y, "근거/출처: " + prov)
        y -= 0.6 * cm

    if extra_notes:
        c.setFont("Helvetica", 9)
        c.drawString(x0, y, _truncate(extra_notes, 200))
        y -= 0.6 * cm

    # Table (top rows)
    tbl_df = result_df.copy()
    if len(tbl_df) > 20:
        tbl_df = tbl_df.head(20)
    if tbl_df.shape[1] > 8:
        tbl_df = tbl_df.iloc[:, :8]

    data = [list(map(str, tbl_df.columns.tolist()))] + [list(map(lambda v: _truncate(str(v), 40), row)) for row in tbl_df.values.tolist()]
    table = Table(data, colWidths=None)

    style = TableStyle([
        ("FONT", (0,0), (-1,0), "Helvetica-Bold", 8),
        ("FONT", (0,1), (-1,-1), "Helvetica", 8),
        ("BACKGROUND", (0,0), (-1,0), colors.lightgrey),
        ("GRID", (0,0), (-1,-1), 0.25, colors.grey),
        ("VALIGN", (0,0), (-1,-1), "TOP"),
        ("LEFTPADDING", (0,0), (-1,-1), 3),
        ("RIGHTPADDING", (0,0), (-1,-1), 3),
        ("TOPPADDING", (0,0), (-1,-1), 2),
        ("BOTTOMPADDING", (0,0), (-1,-1), 2),
    ])
    table.setStyle(style)

    # place table
    tw, th = table.wrapOn(c, width - 3*cm, y - 2*cm)
    # Reserve space for chart if present
    chart_h = 7.0 * cm if fig_png else 0
    max_table_h = max(3*cm, y - (2.0*cm + chart_h))
    if th > max_table_h:
        # shrink by taking fewer rows
        max_rows = max(6, int((max_table_h / (0.45*cm)))-1)
        tbl_df2 = result_df.head(max_rows)
        if tbl_df2.shape[1] > 8:
            tbl_df2 = tbl_df2.iloc[:, :8]
        data = [list(map(str, tbl_df2.columns.tolist()))] + [list(map(lambda v: _truncate(str(v), 40), row)) for row in tbl_df2.values.tolist()]
        table = Table(data)
        table.setStyle(style)
        tw, th = table.wrapOn(c, width - 3*cm, y - 2*cm)

    table.drawOn(c, x0, y - th)
    y = y - th - 0.6*cm

    # Chart image
    if fig_png and y > 2.0*cm:
        try:
            img = ImageReader(BytesIO(fig_png))
            img_w = width - 3*cm
            img_h = min(7.0*cm, y - 1.5*cm)
            c.drawImage(img, x0, y - img_h, width=img_w, height=img_h, preserveAspectRatio=True, anchor="nw")
        except Exception:
            pass

    c.showPage()
    c.save()
    buf.seek(0)
    return buf.getvalue()

def run_analysis(llm: ChatOpenAI, persona: str, selected_files: List[CachedFile]):
    st.subheader("📊 분석 모드 (표/그래프 생성)")
    st.caption("※ 현재 버전은 '엑셀/CSV'를 주요 데이터로 사용하고, PDF/PPT/이미지는 근거/설명 참고용으로 활용합니다.")

    # pick dataset (excel/csv only)
    table_files = [f for f in selected_files if f.kind == "table" and f.df is not None]
    if not table_files:
        st.warning("분석 모드는 현재 엑셀/CSV 파일이 필요합니다. (PDF/PPT/이미지는 근거/설명용)")
        return

    dataset_name = st.selectbox("분석에 사용할 데이터(엑셀/CSV)를 선택하세요", [f.name for f in table_files], index=0)
    df = next(f.df for f in table_files if f.name == dataset_name)
    if df is None or df.empty:
        st.warning("선택한 데이터가 비어 있습니다.")
        return

    # prevent runaway
    if len(df) > MAX_ANALYSIS_ROWS:
        st.info(f"데모 안정성을 위해 상위 {MAX_ANALYSIS_ROWS:,}행만 분석합니다. (원본: {len(df):,}행)")
        df = df.head(MAX_ANALYSIS_ROWS)

    # Quick profile for user (not required)
    with st.expander("데이터 미리보기/컬럼 정보", expanded=False):
        st.write(f"- rows: {len(df):,}, cols: {len(df.columns):,}")
        st.write("- columns:", list(map(str, df.columns.tolist())))
        st.dataframe(df.head(30))

    # Optional: extract/apply rules from docs (PDF/PPT/DOCX)
    doc_files = [f for f in selected_files if f.kind in {"pdf", "ppt", "docx"}]
    apply_doc_rules = st.checkbox(
        "문서 규칙(단가/기간/조건) 추출·적용",
        value=False,
        help="PDF/PPT/DOCX에서 기간/조건/단가 규칙을 추출해 엑셀 분석에 반영합니다. (룰 기반, 데모 안정형)"
    )
    selected_rule_docs: List[CachedFile] = []
    selected_rule_doc_names: List[str] = []
    if apply_doc_rules:
        if not doc_files:
            st.info("선택된 파일 중 PDF/PPT/DOCX가 없습니다. (규칙 추출을 위해 문서 파일을 선택해 주세요)")
        else:
            selected_rule_doc_names = st.multiselect(
                "규칙 추출에 사용할 문서(PDF/PPT/DOCX)",
                options=[f.name for f in doc_files],
                default=[f.name for f in doc_files],
                help="기간/조건/단가가 적힌 문서만 선택하세요. (너무 많은 문서는 정확도/속도에 불리)")
            selected_rule_docs = [f for f in doc_files if f.name in selected_rule_doc_names]

    # Analysis prompt input
    user_query = st.text_area(
        "원하는 분석을 적어주세요",
        value="거래처별 정산금액 합계를 비교하는 막대그래프와 상위 20개 표를 만들어줘",
        height=90
    )
    run_btn = st.button("분석 실행", type="primary")

    if not run_btn:
        return

    # Apply doc-derived rules (unit price / date range / conditions) if enabled
    df_for_plan = df
    rules = None
    rules_prov = ""
    rules_source = ""
    if apply_doc_rules and selected_rule_docs:
        with st.spinner("문서 규칙 추출 중..."):
            docs_text = _doc_text_for_rules(selected_rule_docs)
            rp = _rules_prompt(df_for_plan, docs_text, selected_rule_doc_names)
            msgs = [
                SystemMessage(content="You extract business rules from documents. Output JSON only, matching the given schema. No extra text."),
                HumanMessage(content=rp),
            ]
            raw = llm.invoke(msgs).content
            rules = _safe_json_load(raw)

        if isinstance(rules, dict):
            df_for_plan, rules_prov = apply_rules_to_df(df_for_plan, rules)
            rules_source = ", ".join(selected_rule_doc_names)
            with st.expander("추출된 문서 규칙(JSON) 보기", expanded=False):
                st.json(rules)
            if rules_prov:
                st.caption(f"적용된 규칙: {rules_prov}")
        else:
            st.warning("문서 규칙을 JSON으로 추출하지 못했습니다. (규칙 없이 분석을 진행합니다.)")

    with st.spinner("분석 계획 생성 중..."):
        system = build_system_prompt(persona, mode="analysis")
        plan_prompt = _analysis_prompt(df_for_plan, dataset_name, user_query)
        messages = [
            SystemMessage(content=system),
            HumanMessage(content=plan_prompt)
        ]
        resp = llm.invoke(messages)
        plan = _safe_json_load(resp.content)

    if plan is None:
        st.error("분석 계획(JSON) 파싱에 실패했습니다. (모델 출력이 JSON 형식이 아닐 수 있습니다)")
        with st.expander("원본 출력 보기"):
            st.code(resp.content)
        return

    # Validate and fix LLM-generated plan (prevent wrong filters/top_n)
    if isinstance(plan, dict):
        plan = validate_analysis_plan(df_for_plan, plan, user_query)

    # force dataset file
    plan["dataset_file"] = dataset_name

    # Execute
    with st.spinner("pandas로 집계/그래프 생성 중..."):
        try:
            result_df, fig, provenance, conversion_failures = execute_analysis_plan(df_for_plan, plan)
        except Exception as e:
            st.error(f"분석 실행 오류: {e}")
            with st.expander("계획(JSON) 보기"):
                st.json(plan)
            return

    # Render outputs
    st.success("완료")
    full_prov = provenance
    
    # 변환 실패 정보 표시 (조용히 무시 금지)
    if conversion_failures:
        with st.expander("⚠️ 숫자 변환 실패 정보", expanded=False):
            for col, info in conversion_failures.items():
                st.write(f"**{col}** ({info.get('agg', 'N/A')}): {info['count']}개 실패")
                if info.get('samples'):
                    st.caption(f"샘플: {', '.join(info['samples'][:3])}")
    if rules_source:
        full_prov += f" | [규칙 근거: {rules_source}]"
    if rules_prov:
        full_prov += f" | 규칙 적용: {rules_prov}"
    st.markdown(f"**근거:** {full_prov}")

    # Show plan for transparency
    with st.expander("분석 계획(JSON) 보기", expanded=False):
        st.json(plan)

    # Table
    st.subheader("비교/집계 표")
    st.dataframe(result_df)

    # Chart
    if fig is not None:
        st.subheader("그래프")
        st.pyplot(fig)
    else:
        st.info("요청한 차트 정보를 만들 수 없었습니다. (x/y 컬럼 확인 필요)")

    st.divider()
    st.subheader("📄 1페이지 PDF 리포트")
    st.caption("표(상위 일부) + 그래프(있다면) + 근거를 1페이지 PDF로 내보냅니다.")

    report_title = st.text_input("리포트 제목", value=f"분석 결과 - {dataset_name}")
    report_subtitle = st.text_input("리포트 부제", value=_truncate(user_query, 120))

    fig_png = _fig_to_png_bytes(fig)
    pdf_bytes = build_pdf_report_bytes(
        title=report_title,
        subtitle=report_subtitle,
        provenance=full_prov,
        result_df=result_df,
        fig_png=fig_png,
        extra_notes=(f"규칙 근거: {rules_source}" if rules_source else "")
    )
    st.download_button(
        "PDF 다운로드",
        data=pdf_bytes,
        file_name="analysis_report.pdf",
        mime="application/pdf"
    )

def main():
    st.set_page_config(page_title="PoV 문서 챗봇 + 분석", page_icon="📦", layout="wide")
    init_state()

    # ChatGPT-style CSS
    st.markdown("""
    <style>
    /* ChatGPT-like chat container */
    .stChatMessage {
        padding: 1.5rem 1rem;
        margin-bottom: 1rem;
        border-radius: 0.5rem;
    }

    /* User message styling */
    .stChatMessage[data-testid="user-message"] {
        background-color: transparent;
    }

    /* Assistant message styling */
    .stChatMessage[data-testid="assistant-message"] {
        background-color: #f7f7f8;
    }

    /* Chat input at bottom */
    .stChatInputContainer {
        position: sticky;
        bottom: 0;
        background-color: white;
        padding: 1rem 0;
        border-top: 1px solid #e5e5e5;
    }

    /* Smooth auto-scroll behavior */
    section[data-testid="stChatMessageContainer"] {
        scroll-behavior: smooth;
    }

    /* Message text styling */
    .stChatMessage p {
        margin-bottom: 0.5rem;
        line-height: 1.6;
    }

    /* Code blocks in messages */
    .stChatMessage code {
        background-color: #f0f0f0;
        padding: 0.2rem 0.4rem;
        border-radius: 0.25rem;
        font-size: 0.9em;
    }

    /* Tables in chat */
    .stChatMessage table {
        margin: 1rem 0;
        border-collapse: collapse;
        width: 100%;
    }

    .stChatMessage table th,
    .stChatMessage table td {
        padding: 0.5rem;
        border: 1px solid #e5e5e5;
    }

    /* Hide default streamlit padding for cleaner look */
    .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
    }

    /* Chat messages container max width for better readability */
    section[data-testid="stChatMessageContainer"] > div {
        max-width: 900px;
        margin: 0 auto;
    }
    </style>
    """, unsafe_allow_html=True)

    st.title("📦 물류/정산 문서 PoV 챗봇 + 분석(표/그래프)")
    st.caption("검색/요약(Q&A) + 분석(집계/그래프)을 한 앱에서")

    api_key, temperature, persona, mode, use_web = sidebar()

    if not api_key:
        st.warning("사이드바에 OpenAI API Key를 입력해 주세요.")
        st.stop()

    llm = ChatOpenAI(model=MODEL_NAME, temperature=temperature, api_key=api_key)

    # Get selected files from cache (all files are selected by default)
    file_names = list(st.session_state.file_cache.keys())
    selected_files: List[CachedFile] = [st.session_state.file_cache[n] for n in file_names]

    # Q&A mode: ChatGPT-like full-screen layout
    if mode.startswith("Q&A"):
        run_qa(llm, persona, use_web, selected_files)

    # Analysis mode: Show file selector
    else:
        if file_names:
            selected_files_names = st.multiselect(
                "이번 작업에 사용할 파일을 선택하세요 (권장: 1~4개).",
                options=file_names,
                default=file_names[: min(4, len(file_names))],
            )
            selected_files = [st.session_state.file_cache[n] for n in selected_files_names if n in st.session_state.file_cache]

        run_analysis(llm, persona, selected_files)

if __name__ == "__main__":
    main()
